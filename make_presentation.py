"""
Gringotts Engineering Presentation — PowerPoint Generator
Requires: pip install python-pptx
Run:      python3 make_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ─── Colors ────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0d, 0x11, 0x17)
SURF    = RGBColor(0x16, 0x1b, 0x22)
SURF2   = RGBColor(0x21, 0x26, 0x2d)
BORDER  = RGBColor(0x30, 0x36, 0x3d)
TEXT    = RGBColor(0xe6, 0xed, 0xf3)
MUTED   = RGBColor(0x8b, 0x94, 0x9e)
GOLD    = RGBColor(0xf0, 0xa5, 0x00)
BLUE    = RGBColor(0x58, 0xa6, 0xff)
GREEN   = RGBColor(0x3f, 0xb9, 0x50)
RED     = RGBColor(0xff, 0x7b, 0x72)
PURPLE  = RGBColor(0xd2, 0xa8, 0xff)
DANGER  = RGBColor(0xda, 0x36, 0x33)
WHITE   = RGBColor(0xff, 0xff, 0xff)

# ─── Dimensions ────────────────────────────────────────────────────────────
W  = Inches(13.33)
H  = Inches(7.5)

# ─── Helpers ───────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, text, l, t, w, h,
       size=14, bold=False, color=TEXT,
       align=PP_ALIGN.LEFT, italic=False,
       font="Calibri", wrap=True):
    """Add a simple single-paragraph text box."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = font
    return txb


def mtb(slide, lines, l, t, w, h,
        default_size=12, default_color=TEXT,
        default_font="Calibri", wrap=True):
    """
    Multi-paragraph text box.
    lines = list of dicts:
      { text, size, bold, color, italic, align, font, bullet, space_before }
    """
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap

    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = line.get("align", PP_ALIGN.LEFT)
        if line.get("space_before"):
            p.space_before = Pt(line["space_before"])

        run = p.add_run()
        run.text = line.get("text", "")
        run.font.size  = Pt(line.get("size",  default_size))
        run.font.bold  = line.get("bold",  False)
        run.font.italic = line.get("italic", False)
        run.font.color.rgb = line.get("color", default_color)
        run.font.name  = line.get("font",  default_font)
    return txb


def slide_title(slide, text, subtitle=None):
    """Standard content slide title bar."""
    # gold rule line approximated by a colored rectangle
    rect = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.04)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = GOLD
    rect.line.fill.background()

    tb(slide, text,
       l=0.5, t=0.42, w=12.33, h=0.55,
       size=22, bold=True, color=GOLD)

    if subtitle:
        tb(slide, subtitle,
           l=0.5, t=0.92, w=12.33, h=0.35,
           size=11, color=MUTED)


def card_rect(slide, l, t, w, h, fill=SURF, border=BORDER, radius=6):
    """Filled rectangle standing in for a card."""
    shape = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    return shape


def add_table(slide, headers, rows, l, t, w, h,
              header_color=SURF2, row_colors=None,
              font_size=9, col_widths=None):
    """
    Add a table with dark theme.
    headers: list of strings
    rows: list of lists of strings
    col_widths: list of relative widths (sum=1), optional
    """
    ncols  = len(headers)
    nrows  = len(rows) + 1
    table  = slide.shapes.add_table(
        nrows, ncols,
        Inches(l), Inches(t), Inches(w), Inches(h)
    ).table

    # column widths
    if col_widths:
        total = Inches(w)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = int(total * cw)

    def style_cell(cell, text, bold=False, color=TEXT,
                   bg=None, align=PP_ALIGN.LEFT, size=None):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size or font_size)
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.name  = "Calibri"
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURF
        # thin border
        for border_name in ("top", "bottom", "left", "right"):
            border_el = cell._tc.get_or_add_tcPr()

    # Header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        style_cell(cell, hdr, bold=True, color=MUTED,
                   bg=SURF2, align=PP_ALIGN.LEFT)

    # Data rows
    for ri, row in enumerate(rows):
        bg = row_colors[ri] if row_colors and ri < len(row_colors) else SURF
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            col = TEXT
            # color hints
            if isinstance(val, str):
                if val.startswith("✓"):   col = GREEN
                elif val.startswith("✕"): col = RED
                elif val.startswith("HIGH"): col = RED
                elif val.startswith("MED"):  col = GOLD
                elif val.startswith("LOW"):  col = MUTED
            style_cell(cell, str(val), bg=bg, color=col)

    return table


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    s = blank_slide(prs)
    set_bg(s)

    # big gold accent bar
    bar = s.shapes.add_shape(1, Inches(0), Inches(3.3), Inches(0.12), Inches(1.4))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    tb(s, "Engineering Presentation · April 2026",
       l=0.6, t=1.5, w=10, h=0.4, size=12, color=MUTED)
    tb(s, "Gringotts",
       l=0.6, t=2.0, w=10, h=1.2, size=64, bold=True, color=GOLD)
    tb(s, "A Ledger-First Payment Orchestration System\nfor Citation Payment Processing",
       l=0.6, t=3.25, w=10, h=0.8, size=18, color=TEXT)

    meta_lines = [
        {"text": "Scope  —  Payments, Obligations, Audit", "size": 12, "color": MUTED},
        {"text": "Replaces  —  Victor · IMP · Report Manager", "size": 12, "color": MUTED},
        {"text": "Stack  —  Python · FastAPI · PostgreSQL", "size": 12, "color": MUTED},
    ]
    mtb(s, meta_lines, l=0.6, t=4.3, w=8, h=1.0)


def slide_02_agenda(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Agenda")

    # Part 1 card
    card_rect(s, 0.5, 1.3, 5.9, 4.3, fill=SURF)
    tb(s, "Part 1 — Current System", l=0.7, t=1.45, w=5.5, h=0.35,
       size=13, bold=True, color=BLUE)
    items1 = [
        "1.  System overview & service map",
        "2.  Victor — Stripe payment processing",
        "3.  IMP — Incoming mail / check payments",
        "4.  Report Manager — Financial reporting",
        "5.  Key architectural pain points",
    ]
    mtb(s, [{"text": x, "size": 11, "color": TEXT} for x in items1],
        l=0.9, t=1.85, w=5.3, h=2.5)

    # Part 2 card
    card_rect(s, 6.9, 1.3, 5.9, 4.3, fill=SURF, border=GOLD)
    tb(s, "Part 2 — Proposed System", l=7.1, t=1.45, w=5.5, h=0.35,
       size=13, bold=True, color=GOLD)
    items2 = [
        "1.  Gringotts design philosophy",
        "2.  Core data models & state machines",
        "3.  Event-driven payment flows",
        "4.  Ledger as source of truth",
        "5.  Support UI walkthrough",
        "6.  Migration strategy",
    ]
    mtb(s, [{"text": x, "size": 11, "color": TEXT} for x in items2],
        l=7.1, t=1.85, w=5.5, h=2.5)

    # Goal bar
    card_rect(s, 0.5, 5.8, 12.33, 0.9, fill=RGBColor(0x1a, 0x16, 0x00), border=GOLD)
    tb(s, "Goal of this session:  Align the team on why the current payment system needs to be replaced, "
          "what Gringotts solves, and how we will migrate without breaking existing flows.",
       l=0.7, t=5.9, w=12.0, h=0.7, size=11, color=MUTED)


def slide_03_divider(prs, num, title, desc):
    s = blank_slide(prs)
    set_bg(s)

    # big number watermark
    tb(s, num, l=8.5, t=0.5, w=4.5, h=4.0,
       size=200, bold=True, color=RGBColor(0x1a, 0x20, 0x28),
       align=PP_ALIGN.RIGHT)

    # gold accent bar
    bar = s.shapes.add_shape(1, Inches(0.5), Inches(2.0), Inches(0.08), Inches(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    tb(s, title, l=0.75, t=2.1, w=9, h=1.4,
       size=56, bold=True, color=TEXT)
    tb(s, desc,  l=0.75, t=3.65, w=8.5, h=1.2,
       size=14, color=MUTED)


def slide_04_services(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Current System — Services & Payment Paths")

    cards = [
        ("Victor",          "Online + In-Person",
         "FastAPI. Stripe webhooks → updates Citations directly.\nLEO marks in-person via PUT /mark-paid-in-person.\nStack: FastAPI · Stripe SDK · SQLAlchemy",
         RED),
        ("IMP",             "Physical Checks",
         "Single Flask file (1,124 lines). 8 workflows triggered via HTTP\nwhen scanned mail arrives. Marks citations paid via check.\nStack: Flask · boto3 · S3",
         RED),
        ("Report Manager",  "Financial Reporting",
         "Multi-process daemon. Queries Stripe + Checkbook + Citations\ntables every 4 hrs. Generates CSV → S3.\nStack: FastAPI · Pandas · S3",
         BLUE),
    ]

    for i, (name, sub, desc, col) in enumerate(cards):
        x = 0.5 + i * 4.28
        card_rect(s, x, 1.1, 4.0, 2.2, fill=SURF, border=col)
        tb(s, name, l=x+0.15, t=1.2,  w=3.7, h=0.35, size=13, bold=True, color=col)
        tb(s, sub,  l=x+0.15, t=1.55, w=3.7, h=0.3,  size=11, color=TEXT)
        tb(s, desc, l=x+0.15, t=1.9,  w=3.7, h=1.2,  size=9,  color=MUTED)

    # Architecture flow
    card_rect(s, 0.5, 3.5, 12.33, 3.4, fill=SURF2)
    flow_lines = [
        {"text": "Stripe webhook (charge.succeeded)  →  Victor  →  creates StripePaymentCharges + StripeBalanceTransaction", "size": 9.5, "color": TEXT, "font": "Courier New"},
        {"text": "                                         └──→  mutates Citations: citation_amount_paid += amount, payment_status = PAID|PARTIALLY_PAID", "size": 9.5, "color": MUTED, "font": "Courier New"},
        {"text": "", "size": 5, "color": TEXT},
        {"text": "LEO Dashboard (in-person)  →  Victor  →  mutates Citations: amount_paid = due_amount, source = \"in_person\"", "size": 9.5, "color": TEXT, "font": "Courier New"},
        {"text": "                                  └──→  NO financial record created. No audit log.", "size": 9.5, "color": RED, "font": "Courier New", "bold": True},
        {"text": "", "size": 5, "color": TEXT},
        {"text": "Scanned mail (HTTP POST + S3)  →  IMP  →  creates CheckbookCheckDetails(amount=4000  HARDCODED)", "size": 9.5, "color": TEXT, "font": "Courier New"},
        {"text": "                                      └──→  mutates Citations: stage=payment, status=complete, source=\"check\"", "size": 9.5, "color": MUTED, "font": "Courier New"},
        {"text": "", "size": 5, "color": TEXT},
        {"text": "Cron every 4 hrs  →  Report Manager  →  joins StripePaymentCharges + CheckbookCheckDetails + Citations  →  CSV → S3", "size": 9.5, "color": BLUE, "font": "Courier New"},
    ]
    mtb(s, flow_lines, l=0.7, t=3.6, w=12.0, h=3.2)


def slide_05_how_payments(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Current System — How Payments Work Today")

    cols = [
        {
            "title": "Online (Stripe)",
            "steps": [
                ("charge.succeeded webhook", BLUE),
                ("↓ validate sig + dedupe on stripe_event_id", MUTED),
                ("Create StripePaymentCharge", TEXT),
                ("↓ amount_paid += captured - convenience_fee", MUTED),
                ("Mutate Citations directly", RED),
                ("↓ amount_paid ≥ due_amount?", MUTED),
                ("PAID / PARTIALLY_PAID", GREEN),
            ],
            "good": ["✓ Deduplicates on stripe_event_id", "✓ Saves raw webhook to S3"],
            "bad":  [],
        },
        {
            "title": "In-Person (Victor)",
            "steps": [
                ("PUT /mark-paid-in-person", BLUE),
                ("↓ validate user.customer == citation.customer", MUTED),
                ("Mutate Citations directly", RED),
                ("amount_paid = due_amount ✓", GREEN),
                ("payment_source = \"in_person\"", RED),
            ],
            "good": [],
            "bad":  ["✕ No StripePaymentCharges record",
                     "✕ No CitationActivityLog",
                     "✕ Invisible to reporting"],
        },
        {
            "title": "Check (IMP)",
            "steps": [
                ("POST /add-mail-with-urls", BLUE),
                ("↓ validate citation exists", MUTED),
                ("CheckbookCheckDetails(amount=4000)", RED),
                ("↓ deposited_at = cleared_at = now()", MUTED),
                ("Mutate Citations directly", RED),
                ("CitationActivityLog created ✓", GREEN),
            ],
            "good": [],
            "bad":  ["✕ Amount always $40 — hardcoded",
                     "✕ Deposited = Cleared = now()",
                     "✕ No rollback if check returned"],
        },
    ]

    for i, col in enumerate(cols):
        x = 0.4 + i * 4.28
        tb(s, col["title"], l=x, t=1.1, w=4.0, h=0.35,
           size=13, bold=True, color=BLUE)

        card_rect(s, x, 1.5, 4.0, 3.3, fill=SURF)
        step_lines = []
        for step_text, step_color in col["steps"]:
            step_lines.append({"text": step_text, "size": 9.5, "color": step_color, "font": "Calibri"})
        mtb(s, step_lines, l=x+0.1, t=1.55, w=3.8, h=3.1)

        if col["good"] or col["bad"]:
            card_rect(s, x, 4.95, 4.0, 1.8,
                      fill=SURF if col["good"] else RGBColor(0x1f, 0x10, 0x0d),
                      border=GREEN if col["good"] else RED)
            fb_lines = []
            for g in col["good"]:
                fb_lines.append({"text": g, "size": 10, "color": GREEN})
            for b in col["bad"]:
                fb_lines.append({"text": b, "size": 10, "color": RED})
            mtb(s, fb_lines, l=x+0.15, t=5.05, w=3.7, h=1.6)


def slide_06_whats_broken(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Current System — What's Broken")

    sections = [
        ("Financial Correctness", RED, [
            ("No immutable ledger",
             "All balances live in mutable Citations fields. Refunds, chargebacks, and returned\nchecks don't reverse citation_amount_paid — the balance is permanently wrong."),
            ("Hardcoded check amount ($40) + wrong timestamps",
             "IMP always records $40 regardless of the actual check. check_deposited_at and\ncheck_cleared_at are both set to the moment IMP runs — checks take days to clear."),
            ("In-person leaves no financial record",
             "Victor mutates the citation but creates no payment record.\nIn-person revenue is invisible to Report Manager."),
        ]),
        ("Reliability", RED, [
            ("No transactional safety",
             "Victor's 3 DB writes (charge, citation, balance transaction) are not in a single\ntransaction. A mid-flight crash leaves the DB in a partial state."),
            ("Duplicate payment race condition",
             "Victor and IMP share no lock on citations. Both can mark the same citation paid simultaneously."),
        ]),
        ("Auditability & Reporting", RED, [
            ("Audit trail is incomplete and inconsistent",
             "IMP creates CitationActivityLog for most workflows. Victor creates none —\nonline and in-person payment history is unrecoverable from the DB alone."),
            ("Reports join drifting tables",
             "Report Manager joins StripePaymentCharges, CheckbookCheckDetails, and Citations\n— sources that can fall out of sync."),
        ]),
        ("Maintainability", RED, [
            ("No state machine — string array validation",
             "IMP validates citation state with per-workflow hardcoded lists of valid (stage, status) tuples.\nNo shared enforcement — easy to add transitions that bypass business rules."),
            ("Business logic scattered across services",
             "Payment rules live in Victor's webhook controller, IMP's Flask routes,\nand Report Manager's SQL — no single place to reason about what a payment means."),
        ]),
    ]

    positions = [(0.4, 1.1), (6.9, 1.1), (0.4, 4.0), (6.9, 4.0)]

    for idx, ((title, col, items), (x, y)) in enumerate(zip(sections, positions)):
        card_rect(s, x, y, 6.0, 2.85, fill=RGBColor(0x1a, 0x10, 0x0d), border=RED)
        tb(s, title, l=x+0.15, t=y+0.1, w=5.7, h=0.3,
           size=11, bold=True, color=col)
        lines = []
        for i, (item_title, item_desc) in enumerate(items):
            if i > 0:
                lines.append({"text": "", "size": 4, "color": TEXT})
            lines.append({"text": "✕  " + item_title, "size": 10, "bold": True, "color": RED})
            lines.append({"text": item_desc, "size": 9, "color": MUTED})
        mtb(s, lines, l=x+0.15, t=y+0.45, w=5.7, h=2.3)


def slide_07_gringotts_divider(prs):
    slide_03_divider(prs, "02", "Gringotts",
                     "A ledger-first, event-driven payment orchestration system built on\nimmutability, auditability, and clean separation of concerns.")


def slide_08_philosophy(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Gringotts — Design Philosophy")

    principles = [
        ("Ledger is Source of Truth",
         "All balances are derived from immutable ledger entries.\nCitation fields and cart statuses are computed, never authoritative.", GOLD),
        ("Immutability Over Mutation",
         "Financial history is append-only. No overwriting. A full payment\ntimeline is reconstructable from the ledger alone.", GOLD),
        ("Event-Driven State",
         "Vendor signals (Stripe webhooks, CheckAlt drops) drive state transitions\nasynchronously. The API never waits for vendors.", BLUE),
        ("Idempotency by Default",
         "Every processor is safe to replay. Duplicate webhook delivery, queue\nretries, or manual replays produce the same result.", BLUE),
        ("Vendor Agnostic",
         "Core event processor has zero vendor imports. Adding a new payment\nprocessor means adding one vendor module.", GREEN),
        ("Full Auditability",
         "Every state change produces an AuditLog with before/after snapshots\nand correlation IDs linking related events.", GREEN),
    ]

    for i, (title, desc, col) in enumerate(principles):
        col_i = i % 3
        row_i = i // 3
        x = 0.4 + col_i * 4.28
        y = 1.2 + row_i * 2.2
        card_rect(s, x, y, 4.0, 2.0, fill=SURF)
        tb(s, title, l=x+0.15, t=y+0.1, w=3.7, h=0.35,
           size=12, bold=True, color=col)
        tb(s, desc, l=x+0.15, t=y+0.5, w=3.7, h=1.4,
           size=10, color=MUTED)

    # Layer separation
    card_rect(s, 0.4, 5.6, 12.33, 1.25, fill=RGBColor(0x0a, 0x14, 0x0a))
    flow_parts = [
        ("Intent Layer", "API + Cart", GOLD),
        ("→", "", MUTED),
        ("Processing Layer", "Events + Ledger", BLUE),
        ("→", "", MUTED),
        ("Derived State", "Carts + Obligations + Citations", GREEN),
    ]
    x = 1.0
    for name, sub, col in flow_parts:
        if name == "→":
            tb(s, "→", l=x, t=5.9, w=0.4, h=0.5, size=18, color=MUTED, align=PP_ALIGN.CENTER)
            x += 0.5
        else:
            tb(s, name, l=x, t=5.68, w=3.5, h=0.35, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
            tb(s, sub,  l=x, t=6.05, w=3.5, h=0.3,  size=9,  color=MUTED, align=PP_ALIGN.CENTER)
            x += 3.6


def slide_09_data_models(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Core Data Models")

    models = [
        ("Obligation", GOLD,
         "What someone owes. One per citation fee, late fee,\nflagging fee, or convenience fee.",
         ["amount: original owed",
          "allocated_total: sum of settled ledger allocations",
          "outstanding_amount = amount − allocated_total",
          "overpaid_amount: if allocated > amount",
          "waived_amount: forgiven (not paid)",
          "status: new → open ↔ locked → closed\n        (↓ waived, ↓ superseded, ↓ disputed)"]),
        ("Cart", BLUE,
         "Groups obligations for a single payment intent.\nOne Stripe charge per cart — not per citation.",
         ["draft → checkout → payment_submitted",
          "→ payment_confirmed → payment_settled",
          "(↓ abandoned, ↓ disputed → won/lost)"]),
        ("Ledger", GREEN,
         "Append-only. Every payment event recorded here\nwith allocations to specific obligations.",
         ["stage: payment_confirmed / settled / refunded / dispute_*",
          "allocations: [{cart_item, obligation_id, amount}]",
          "Each obligation's balance = SUM(its allocations)"]),
        ("VendorEvent", PURPLE,
         "Normalized inbound signal from Stripe/CheckAlt.\nDrives all state transitions.",
         ["event_type: payment_confirmed / settled / refunded / dispute_*",
          "transaction: normalized payment data",
          "payment_breakdown: [{type, amount, label}]",
          "item_allocations: [{owner_id, amount}]"]),
    ]

    for i, (name, col, desc, fields) in enumerate(models):
        col_i = i % 2
        row_i = i // 2
        x = 0.4 + col_i * 6.55
        y = 1.1 + row_i * 2.95
        card_rect(s, x, y, 6.2, 2.75, fill=SURF)

        # badge
        badge = s.shapes.add_shape(1, Inches(x+0.15), Inches(y+0.12),
                                   Inches(1.4), Inches(0.28))
        badge.fill.solid(); badge.fill.fore_color.rgb = col
        badge.line.fill.background()
        tb(s, name, l=x+0.18, t=y+0.12, w=1.36, h=0.28,
           size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)

        tb(s, desc,  l=x+0.15, t=y+0.47, w=5.9, h=0.55, size=10, color=MUTED)
        field_lines = [{"text": f, "size": 9, "color": TEXT, "font": "Courier New"} for f in fields]
        mtb(s, field_lines, l=x+0.15, t=y+1.05, w=5.9, h=1.55)


def slide_10_stripe_flow(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Online Payment Flow (Stripe)")

    # Synchronous column
    card_rect(s, 0.4, 1.1, 5.5, 3.4, fill=SURF)
    tb(s, "Synchronous (User-Facing API)", l=0.55, t=1.15,
       w=5.2, h=0.3, size=11, bold=True, color=BLUE)
    sync_lines = [
        {"text": "POST /checkout  →  Create cart (draft → checkout)", "size": 10, "color": TEXT, "font": "Courier New"},
        {"text": "  Lock selected obligations. Return payment intent.", "size": 9, "color": MUTED, "font": "Courier New"},
        {"text": "", "size": 5, "color": TEXT},
        {"text": "POST /submit  →  202 Accepted", "size": 10, "color": GREEN, "font": "Courier New"},
        {"text": "  Cart → payment_submitted. Obligations locked.", "size": 9, "color": MUTED, "font": "Courier New"},
        {"text": "  Returns IMMEDIATELY — no waiting for Stripe.", "size": 9, "bold": True, "color": GOLD, "font": "Courier New"},
    ]
    mtb(s, sync_lines, l=0.55, t=1.5, w=5.2, h=2.8)

    # Asynchronous column
    card_rect(s, 6.2, 1.1, 6.63, 3.4, fill=SURF)
    tb(s, "Asynchronous (Event Worker)", l=6.35, t=1.15,
       w=6.3, h=0.3, size=11, bold=True, color=GOLD)
    async_lines = [
        {"text": "① Stripe webhook received — validated + deduplicated on stripe_event_id", "size": 9.5, "color": TEXT},
        {"text": "② Transaction upsert — normalize Stripe charge into vendor-agnostic Transaction", "size": 9.5, "color": TEXT},
        {"text": "③ VendorEvent created — type: payment_confirmed. Enqueued for processing.", "size": 9.5, "color": TEXT},
        {"text": "④ EventProcessor runs (2-pass):", "size": 9.5, "color": TEXT},
        {"text": "    Pass 1: Resolve allocations to obligations", "size": 9, "color": MUTED},
        {"text": "    Pass 2: Attach cart items → record in ledger", "size": 9, "color": MUTED},
        {"text": "⑤ Derived state updated — Cart → payment_confirmed.", "size": 9.5, "color": TEXT},
        {"text": "    Obligation balances recalculated from ledger.", "size": 9, "color": MUTED},
    ]
    mtb(s, async_lines, l=6.35, t=1.5, w=6.3, h=2.8)

    # Key differences
    card_rect(s, 0.4, 4.7, 12.33, 2.5, fill=RGBColor(0x0a, 0x14, 0x0a), border=GREEN)
    tb(s, "Key Differences from Victor", l=0.55, t=4.77, w=12.0, h=0.28,
       size=11, bold=True, color=GREEN)
    diff_lines = [
        {"text": "✓  API never blocks on vendor — Submit returns 202 immediately. Vendor confirmation is async. UI polls for status.", "size": 9.5, "color": GREEN},
        {"text": "✓  One charge per cart (not per citation) — Violator pays all citations in one Stripe charge. Fewer transactions = lower fees.", "size": 9.5, "color": GREEN},
        {"text": "✓  Ledger is updated, not Citations table — Obligation balances derived from ledger allocations. No more mutable citation_amount_paid.", "size": 9.5, "color": GREEN},
        {"text": "✓  Full idempotency — Processing the same webhook twice yields the same ledger state. Safe to replay.", "size": 9.5, "color": GREEN},
    ]
    mtb(s, diff_lines, l=0.55, t=5.1, w=12.1, h=2.0)


def slide_11_check_flow(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Check & In-Person Payment Flow (Proposed)")

    tb(s, "Vendor events carry the actual check amount — no hardcoding. Future integration with CheckAlt or any mail processor uses the same pipeline.",
       l=0.5, t=1.0, w=12.33, h=0.4, size=10, color=MUTED)

    steps = [
        ("1", "Check vendor event ingested",
         "Contains citation reference, actual amount, check metadata.\nVendor module normalizes into a VendorEvent."),
        ("2", "VendorEvent created (no pre-existing cart)",
         "Allocations are payable type: (owner_id, owner_type, label)\n— resolved dynamically."),
        ("3", "EventProcessor: Pass 1",
         "Resolves payable allocations → finds matching open obligation.\nCreates system_created cart on-the-fly."),
        ("4", "Ledger records actual payment amount",
         "Deposited and cleared timestamps tracked separately.\nObligation balance recalculated from ledger."),
        ("5", "Partial payment? Supported natively.",
         "Multiple checks can partially pay one obligation.\nEach adds a ledger allocation. Obligation closes when outstanding = 0."),
    ]

    for i, (num, title, desc) in enumerate(steps):
        x = 0.4 + i * 2.5
        card_rect(s, x, 1.55, 2.3, 1.85, fill=SURF)
        tb(s, num, l=x+0.1, t=1.62, w=0.3, h=0.3,
           size=13, bold=True, color=GOLD)
        tb(s, title, l=x+0.1, t=1.97, w=2.1, h=0.4,
           size=9.5, bold=True, color=TEXT)
        tb(s, desc, l=x+0.1, t=2.4, w=2.1, h=0.9,
           size=8.5, color=MUTED)

    # Comparison table
    tb(s, "IMP vs Gringotts: Check & In-Person", l=0.4, t=3.6, w=12.33, h=0.3,
       size=11, bold=True, color=BLUE)

    headers = ["Aspect", "Today (IMP / Victor)", "Gringotts"]
    rows = [
        ["Check amount",         "Always $40 (hardcoded in IMP)",         "Actual amount from vendor event"],
        ["Partial payments",      "Not supported",                          "Supported natively"],
        ["Returned checks",       "No financial rollback",                  "Refund event reverses ledger allocation"],
        ["Audit trail",           "CitationActivityLog (partial)",          "Full ledger + AuditLog with before/after"],
        ["In-person record",      "None — invisible to reporting",          "VendorEvent + ledger entry; full audit trail"],
        ["Duplicate safety",      "No deduplication key (IMP)",             "Deduped on vendor event ID"],
        ["State management",      "String array validation",                "Deterministic state machine"],
    ]
    add_table(s, headers, rows, l=0.4, t=3.95, w=12.33, h=3.25,
              font_size=9, col_widths=[0.22, 0.39, 0.39])


def slide_12_ledger_truth(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Ledger as Source of Truth")

    # Code block - balance computation
    card_rect(s, 0.4, 1.1, 7.5, 3.6, fill=SURF2)
    tb(s, "How obligation balances are computed:", l=0.55, t=1.15,
       w=7.2, h=0.28, size=10, bold=True, color=BLUE)
    code = """def update_obligation_according_to_ledger(obligation):
    # Recalculate from scratch — every time
    allocations = ledger.get_allocations_for(obligation)

    allocated_total   = sum(a.amount for a in allocations)
    outstanding_amount = max(obligation.amount - allocated_total, 0)
    overpaid_amount    = max(allocated_total - obligation.amount, 0)

    if outstanding_amount == 0:
        obligation.status = CLOSED
    elif obligation.status == LOCKED:
        obligation.status = OPEN   # reopen if refunded

    obligation.allocated_total    = allocated_total
    obligation.outstanding_amount = outstanding_amount"""
    tb(s, code, l=0.55, t=1.48, w=7.2, h=3.1,
       size=9, color=TEXT, font="Courier New")

    tb(s, "Key insight: If you delete all ledger entries for an obligation, the balance resets to zero.\nYou can replay any event and get the correct answer. State cannot drift.",
       l=0.4, t=4.85, w=7.5, h=0.55, size=10, bold=True, color=GOLD)

    # Right column
    card_rect(s, 8.1, 1.1, 4.8, 2.3, fill=SURF)
    tb(s, "Ledger entry lifecycle:", l=8.25, t=1.15, w=4.5, h=0.28, size=10, bold=True, color=TEXT)
    lifecycle_lines = [
        {"text": "payment_confirmed — Vendor confirmed receipt", "size": 9.5, "color": BLUE},
        {"text": "↓ (charge clears)", "size": 9, "color": MUTED},
        {"text": "payment_settled — Funds cleared (stage update only)", "size": 9.5, "color": GREEN},
        {"text": "↓ (refund requested)", "size": 9, "color": MUTED},
        {"text": "payment_refunded — New entry; allocations reversed", "size": 9.5, "color": RED},
    ]
    mtb(s, lifecycle_lines, l=8.25, t=1.48, w=4.5, h=1.8)

    # Idempotency
    card_rect(s, 8.1, 3.55, 4.8, 2.1, fill=SURF2)
    tb(s, "Idempotency guarantee:", l=8.25, t=3.6, w=4.5, h=0.28, size=10, bold=True, color=TEXT)
    idempotency_code = """existing = ledger.identify_existing_record(
    vendor_name=..., transaction_id=...,
    type=..., label=..., cart_id=...
)
if existing.stage in [CONFIRMED, SETTLED]:
    return existing  # skip — already processed
# Only append if new"""
    tb(s, idempotency_code, l=8.25, t=3.93, w=4.5, h=1.6,
       size=8.5, color=TEXT, font="Courier New")

    card_rect(s, 8.1, 5.8, 4.8, 0.85, fill=RGBColor(0x0a, 0x14, 0x0a), border=GREEN)
    tb(s, "Reports query only settled_at IS NOT NULL ledger entries.\nOne source. No joins across drifting tables. Always consistent.",
       l=8.25, t=5.88, w=4.5, h=0.7, size=9, color=GREEN)


def slide_13_obligation_ops(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Obligation Operations (Support Actions)",
                subtitle="Gringotts supports three first-class operations for adjusting obligations — all with full audit trails.")

    ops = [
        ("Waive", "Forgive Unpaid Debt", BLUE,
         "For unpaid obligations (no prior payments).",
         ["Precondition: allocated_total == 0",
          "Effect: waived_amount += outstanding",
          "         outstanding_amount = 0",
          "         status = WAIVED (terminal)",
          "",
          "Use case: Court grants hardship exemption.",
          "Obligation forgiven without any payment."]),
        ("Partial Waive", "Forgive Remaining Balance", BLUE,
         "For partially-paid obligations.",
         ["Precondition: allocated_total > 0 AND outstanding > 0",
          "Effect: waived_amount += outstanding",
          "         outstanding_amount = 0",
          "         status = CLOSED",
          "",
          "Use case: Paid $500 of $1000.",
          "Court forgives remaining $500."]),
        ("Supersede", "Replace at New Amount", PURPLE,
         "For open obligations with no prior payments.",
         ["Precondition: status == OPEN AND allocated_total == 0",
          "Effect: old obligation → SUPERSEDED",
          "         new obligation: amount = new_amount",
          "                        status = OPEN",
          "",
          "Use case: Court revises fine from $1000 to $2500.",
          "Both old and new in audit trail."]),
    ]

    for i, (name, sub, col, desc, fields) in enumerate(ops):
        x = 0.4 + i * 4.28
        card_rect(s, x, 1.35, 4.0, 4.6, fill=SURF)
        badge = s.shapes.add_shape(1, Inches(x+0.15), Inches(1.45),
                                   Inches(1.6), Inches(0.28))
        badge.fill.solid(); badge.fill.fore_color.rgb = col
        badge.line.fill.background()
        tb(s, name, l=x+0.17, t=1.45, w=1.56, h=0.28,
           size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)
        tb(s, sub, l=x+0.15, t=1.8, w=3.7, h=0.3,
           size=11, bold=True, color=TEXT)
        tb(s, desc, l=x+0.15, t=2.15, w=3.7, h=0.3, size=10, color=MUTED)
        field_lines = [{"text": f, "size": 9.5, "color": TEXT, "font": "Courier New"} for f in fields]
        mtb(s, field_lines, l=x+0.15, t=2.5, w=3.7, h=3.2)

    card_rect(s, 0.4, 6.1, 12.33, 0.9, fill=SURF2)
    tb(s, "All operations produce AuditLog entries with full before/after snapshots and correlation IDs. "
          "Support staff can trace every obligation change to a specific action and actor.",
       l=0.55, t=6.2, w=12.0, h=0.7, size=10, color=MUTED)


def slide_14_architecture(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Gringotts — System Architecture")

    card_rect(s, 0.4, 1.1, 12.33, 5.9, fill=SURF2)
    arch = """\
┌────────────────────────────────────────────────────────────────────┐
│                        USER / SUPPORT LAYER                        │
│  Public API                       Internal Support API             │
│  POST /checkout  → 201            POST /refunds              → 202 │
│  POST /submit    → 202            POST /obligations/{id}/waive     │
│                                   GET  /reports/ledger             │
└──────────────────┬──────────────────────────┬──────────────────────┘
                   │ (sync, returns fast)      │
                   ↓                           ↓
      ┌────────────────────┐       ┌────────────────────────┐
      │  CartProcessor     │       │  RefundProcessor       │
      │  create_cart       │       │  Dequeue refund request│
      │  checkout_cart     │       │  Call VendorAdapter    │
      │  submit → lock     │       │  Emit payment_refunded │
      └──────────┬─────────┘       └────────────────────────┘
                 │
      ┌──────────┴─────────────────────────────────────────────────┐
      │                   VENDOR EVENT QUEUE                       │
      │  (Stripe webhooks · CheckAlt drops · Court API)            │
      │  Deduplicated by vendor_event_id                           │
      └──────────────────────────┬─────────────────────────────────┘
                                 │
                                 ↓
              ┌─────────────────────────────────────┐
              │   EventProcessor (Core Engine)      ├── if anomaly ──→ ┌───────────────────┐
              │   Pass 1: Resolve allocations        │                  │   TRIAGE QUEUE    │
              │     → CART_ITEM:  find CartItem      │                  │   invalid_payment │
              │     → PAYABLE:    match (owner,label)│                 │   system_error    │
              │   Pass 2: Attach cart items          │                  └─────────┬─────────┘
              │   → Ledger.record_event()            │                            ↓
              └────────────┬────────────────────────┘                  ┌───────────────────┐
                           │                                            │  TriageProcessor  │
          ┌────────────────┼────────────────┐                          │  unresolved allocs│
          ↓                ↓                ↓                          │  overpayments     │
   ┌───────────┐   ┌────────────┐   ┌────────────────┐                │  ledger↔OBL delta │
   │  Ledger   │   │   Carts    │   │  Obligations   │                │  → TriageRequest  │
   │ (Append   │   │ (Derived)  │   │  (Derived)     │                └───────────────────┘
   │  Only)    │   └────────────┘   └────────────────┘
   └───────────┘
          │
   ┌──────┴────────────────────────┐
   │   AuditLog + AccessLog        │
   │   correlation_id chaining     │
   └───────────────────────────────┘"""
    tb(s, arch, l=0.55, t=1.2, w=12.1, h=5.7,
       size=8, color=TEXT, font="Courier New")


def slide_15_ui_divider(prs):
    slide_03_divider(prs, "03", "Support UI",
                     "How support staff will interact with the Gringotts system —\ncitation lookup, payment history, refunds, obligation adjustments, and audit logs.")


def slide_16_citation_overview(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Support UI — Citation Payment Overview",
                subtitle="Support staff search by citation number. Single pane shows full financial state.")

    # Browser chrome
    card_rect(s, 0.4, 1.2, 12.33, 5.85, fill=SURF2)
    card_rect(s, 0.4, 1.2, 12.33, 0.35, fill=SURF)

    # Nav sidebar
    card_rect(s, 0.4, 1.55, 2.1, 5.5, fill=SURF)
    nav_items = ["⚡ Payment Overview", "📋 Obligation History",
                 "📦 Cart Detail", "📒 Ledger Entries",
                 "🔍 Audit Log", "↩ Refunds"]
    for i, item in enumerate(nav_items):
        col = BLUE if i == 0 else MUTED
        bg = RGBColor(0x1a, 0x24, 0x38) if i == 0 else SURF
        card_rect(s, 0.45, 1.65 + i*0.52, 1.95, 0.45, fill=bg, border=bg)
        tb(s, item, l=0.55, t=1.7 + i*0.52, w=1.8, h=0.35, size=9, color=col)

    tb(s, "ACTIONS", l=0.55, t=4.85, w=1.8, h=0.2, size=7.5, bold=True, color=MUTED)
    actions = ["Waive Obligation", "Supersede Obligation", "Initiate Refund"]
    for i, act in enumerate(actions):
        card_rect(s, 0.5, 5.1 + i*0.38, 1.85, 0.3, fill=SURF2, border=BLUE)
        tb(s, act, l=0.55, t=5.13 + i*0.38, w=1.75, h=0.25, size=8.5, color=BLUE)

    # Main content
    tb(s, "Citation #000322452", l=2.65, t=1.3, w=5, h=0.3, size=13, bold=True, color=TEXT)

    fields = [
        ("Citation Type", "Check payment · partial · OBL-3"),
        ("Total Owed",    "$40.00"),
        ("Outstanding",   "$20.00"),
        ("Status",        "PARTIALLY PAID"),
    ]
    for i, (label, val) in enumerate(fields):
        x = 2.65 + (i % 2) * 4.0
        y = 1.65 + (i // 2) * 0.45
        tb(s, label, l=x, t=y, w=1.5, h=0.25, size=8.5, color=MUTED)
        tb(s, val, l=x+1.6, t=y, w=2.3, h=0.25, size=8.5, color=TEXT)

    tb(s, "Obligations", l=2.65, t=2.6, w=9.8, h=0.28, size=10, bold=True, color=MUTED)
    obl_headers = ["Label", "Amount", "Paid", "Outstanding", "Status", "Action"]
    obl_rows = [
        ["Citation Fee (OBL-3)", "$40.00", "$20.00", "$20.00", "OPEN", "Waive ↗"],
        ["Check (42/1)",          "cart:4", "$20.00 paid", "—", "SETTLED", "—"],
    ]
    add_table(s, obl_headers, obl_rows, l=2.65, t=2.93, w=9.9, h=1.2,
              font_size=9.5, col_widths=[0.25, 0.13, 0.13, 0.16, 0.16, 0.17])

    tb(s, "Key feature: Per-obligation status. Actions available inline based on state. "
          "Outstanding amount always derived from ledger — never stale.",
       l=2.65, t=4.25, w=9.9, h=0.4, size=9, color=MUTED)


def slide_17_sample_data(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Support UI — Obligation & Cart Sample Data")

    # Obligations
    tb(s, "Obligations — Citations 000331714 & 000329181", l=0.4, t=1.05, w=8, h=0.28,
       size=11, bold=True, color=BLUE)
    tb(s, "One obligation per fee type. Balances are always derived from ledger allocations.",
       l=0.4, t=1.35, w=8, h=0.25, size=9, color=MUTED)

    obl_headers = ["id", "label", "amount", "allocated", "outstanding", "status"]
    obl_rows = [
        ["1", "citation_fee",   "$40.00", "$40.00", "$0.00", "CLOSED"],
        ["2", "citation_fee",   "$40.00", "$80.00", "$0.00", "CLOSED (+$40 overpaid)"],
        ["9", "convenience_fee","$4.00",  "$4.00",  "$0.00", "CLOSED"],
    ]
    add_table(s, obl_headers, obl_rows, l=0.4, t=1.65, w=7.8, h=1.3,
              font_size=9.5, col_widths=[0.06, 0.22, 0.14, 0.14, 0.16, 0.28])

    # How outstanding is computed
    card_rect(s, 0.4, 3.1, 7.8, 1.15, fill=SURF2)
    tb(s, "How outstanding_amount is computed:", l=0.55, t=3.15, w=7.5, h=0.25,
       size=9.5, bold=True, color=TEXT)
    compute_lines = [
        {"text": "OBL-1: allocated($40) − amount($40) = $0  → CLOSED", "size": 9.5, "color": GREEN, "font": "Courier New"},
        {"text": "OBL-2: allocated($80) − amount($40) = overpaid $40, outstanding=$0  → CLOSED", "size": 9.5, "color": GOLD, "font": "Courier New"},
        {"text": "OBL-9: allocated($4)  − amount($4)  = $0  → CLOSED", "size": 9.5, "color": GREEN, "font": "Courier New"},
    ]
    mtb(s, compute_lines, l=0.55, t=3.42, w=7.5, h=0.8)

    # Cart
    tb(s, "Cart — 1", l=8.8, t=1.05, w=4.8, h=0.28, size=11, bold=True, color=BLUE)
    tb(s, "One cart per checkout session. Groups all obligations into a single vendor charge.",
       l=8.8, t=1.35, w=4.5, h=0.35, size=9, color=MUTED)

    cart_fields = [
        ("id",             "1"),
        ("status",         "payment_settled"),
        ("payment_mode",   "online"),
        ("vendor",         "stripe"),
        ("amount",         "$84.00"),
        ("refund_amount",  "$40.00"),
        ("vendor_ref",     "txn_3RKstt…S1G9OXSH1"),
        ("created_at",     "2026-03-24 11:56:15Z"),
    ]
    card_rect(s, 8.8, 1.65, 4.3, 2.9, fill=SURF)
    for i, (label, val) in enumerate(cart_fields):
        y = 1.72 + i * 0.33
        tb(s, label, l=8.95, t=y, w=1.4, h=0.28, size=9, color=MUTED)
        tb(s, val,   l=10.4, t=y, w=2.6, h=0.28, size=9, color=TEXT, font="Courier New")

    tb(s, "Cart Items:", l=8.8, t=4.7, w=4.3, h=0.28, size=10, bold=True, color=TEXT)
    ci_headers = ["item", "obligation", "label", "amount"]
    ci_rows = [
        ["CI-1", "OBL-9", "convenience_fee",       "$4.00"],
        ["CI-2", "OBL-1", "citation_fee (000331714)", "$40.00"],
        ["CI-3", "OBL-2", "citation_fee (000329181)", "$40.00"],
    ]
    add_table(s, ci_headers, ci_rows, l=8.8, t=5.0, w=4.3, h=1.0,
              font_size=9, col_widths=[0.12, 0.14, 0.48, 0.26])

    tb(s, "OBL-2 shows $80 allocated — paid by both Stripe and CheckAlt, resulting in $40 overpayment.",
       l=0.4, t=6.5, w=12.33, h=0.5, size=9, color=GOLD)


def slide_18_ledger_transitions(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Support UI — Ledger: Data Transitions",
                subtitle="The ledger is append-only. Each vendor event adds rows — obligation balances recalculated from running total of allocations.")

    states = [
        ("State 1: payment_confirmed  (11:56:29Z)", BLUE,
         "Stripe webhook received. Funds not yet settled.",
         [["LDG-1", "CONFIRMED", "cart_payment", "$84.00"],
          ["  ↳ CI-2", "alloc", "OBL-1 citation_fee", "$40.00"],
          ["  ↳ CI-3", "alloc", "OBL-2 citation_fee", "$40.00"],
          ["  ↳ CI-1", "alloc", "OBL-9 conv_fee",     "$4.00"],
          ["LDG-2", "CONFIRMED", "stripe_fee", "-$1.55"]]),
        ("State 2: payment_settled  (11:56:33Z)", GREEN,
         "Stripe charge clears. Stage updated in-place on existing rows.",
         [["LDG-1", "SETTLED", "cart_payment", "$84.00"],
          ["  ↳ CI-2", "alloc", "OBL-1 citation_fee", "$40.00"],
          ["  ↳ CI-3", "alloc", "OBL-2 citation_fee", "$40.00"],
          ["  ↳ CI-1", "alloc", "OBL-9 conv_fee",     "$4.00"],
          ["LDG-2", "SETTLED", "stripe_fee", "-$1.55"]]),
        ("State 3: payment_refunded  (11:57:21Z)", RED,
         "New ledger item appended — nothing overwritten.",
         [["LDG-1", "SETTLED", "cart_payment (prior)", "$84.00"],
          ["LDG-2", "SETTLED", "stripe_fee (prior)", "-$1.55"],
          ["LDG-7", "REFUNDED", "refund_payment", "-$40.00"],
          ["  ↳ CI-2", "alloc", "OBL-1 citation_fee", "-$40.00"]]),
    ]

    headers = ["ldg_id", "stage", "label", "amount"]

    for i, (title, col, desc, rows) in enumerate(states):
        x = 0.4 + i * 4.28
        tb(s, title, l=x, t=1.08, w=4.0, h=0.28, size=9.5, bold=True, color=col)
        tb(s, desc,  l=x, t=1.38, w=4.0, h=0.3,  size=8.5, color=MUTED)

        # mini table
        card_rect(s, x, 1.73, 4.0, 2.8, fill=SURF)
        add_table(s, headers, rows, l=x+0.05, t=1.78, w=3.9, h=2.65,
                  font_size=8.5, col_widths=[0.2, 0.22, 0.38, 0.2])

    # Key insight bar
    card_rect(s, 0.4, 6.05, 12.33, 0.85, fill=RGBColor(0x0a, 0x14, 0x0a), border=GREEN)
    tb(s, "Key insight: The ledger is never overwritten. Refunds are new rows. "
          "Obligation.outstanding_amount = SUM of all ledger allocations — always derivable, always correct.",
       l=0.55, t=6.12, w=12.1, h=0.7, size=10, bold=True, color=GREEN)


def slide_19_access_audit(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Support UI — Access Log & Audit Log")

    # Access log
    tb(s, "Access Log", l=0.4, t=1.08, w=5.9, h=0.3, size=12, bold=True, color=BLUE)
    tb(s, "Every HTTP request and worker execution. Operational observability — latency, errors, actors.",
       l=0.4, t=1.4, w=5.9, h=0.28, size=9, color=MUTED)
    card_rect(s, 0.4, 1.72, 5.9, 3.95, fill=SURF2)
    access_code = """── request ─────────────────────────────────────
channel:        http
request_id:     req_submit_1_1774353383900
correlation_id: cart_payment:1
actor_type:     user / actor_id: demo_user
method:         POST /payments/submit
status_code:    202 / duration_ms: 1
timestamp:      2026-03-24T11:56:23Z

── worker ──────────────────────────────────────
channel:        worker
request_id:     worker_event:1
correlation_id: cart_payment:1
causation_id:   vendor_event:1
method:         process_payment_event
status_code:    200 / duration_ms: 94
timestamp:      2026-03-24T11:56:29Z"""
    tb(s, access_code, l=0.55, t=1.8, w=5.6, h=3.8,
       size=8.5, color=TEXT, font="Courier New")

    # Audit log
    tb(s, "Audit Log", l=6.6, t=1.08, w=6.4, h=0.3, size=12, bold=True, color=GOLD)
    tb(s, "Every domain-level state change with before/after snapshots. Full traceability for any payment complaint.",
       l=6.6, t=1.4, w=6.1, h=0.28, size=9, color=MUTED)
    card_rect(s, 6.6, 1.72, 6.4, 3.95, fill=SURF2)
    audit_code = """── obligation_status_changed ─────────────────
entity_type:    obligation
entity_id:      1 (000331714)
old_status:     LOCKED / new_status: CLOSED
before:         {"outstanding": 4000, "allocated": 0}
after:          {"outstanding": 0,    "allocated": 4000}
correlation_id: cart_payment:1
causation_id:   vendor_event:1
timestamp:      2026-03-24T11:56:29Z

── cart_status_changed ───────────────────────
entity_type:    cart / entity_id: 1
old_status:     payment_confirmed
new_status:     payment_settled
correlation_id: cart_payment:1
causation_id:   vendor_event:2
timestamp:      2026-03-24T11:56:33Z"""
    tb(s, audit_code, l=6.75, t=1.8, w=6.1, h=3.8,
       size=8.5, color=TEXT, font="Courier New")

    # Tracing note
    card_rect(s, 0.4, 5.82, 12.33, 0.85, fill=RGBColor(0x14, 0x14, 0x22), border=PURPLE)
    tb(s, "Tracing a complaint: Citation 000331714 → OBL-1 → Cart 1 → cart_payment:1 → "
          "all access logs + audit logs for the full flow (checkout, submit, confirmed, settled, refund) "
          "across HTTP and worker execution.",
       l=0.55, t=5.9, w=12.1, h=0.7, size=9.5, color=PURPLE)


def slide_20_refund_ui(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Support UI — Initiating a Refund",
                subtitle="Support staff select a settled cart, choose items to refund, and submit. The system handles the rest asynchronously.")

    # Refund form
    card_rect(s, 0.4, 1.25, 5.8, 5.25, fill=SURF)
    tb(s, "Refund Form (Cart 1)", l=0.55, t=1.35, w=5.5, h=0.28,
       size=11, bold=True, color=TEXT)
    tb(s, "Cart Status: payment_settled · Stripe · 2026-03-24 · $84.00",
       l=0.55, t=1.68, w=5.5, h=0.28, size=9, color=MUTED)

    tb(s, "Select Items to Refund:", l=0.55, t=2.05, w=5.5, h=0.25,
       size=9.5, bold=True, color=TEXT)

    items = [
        ("☑", "Citation Fee · OBL-1 (000331714) · CI-2 · $40.00", BLUE),
        ("☐", "Convenience Fee · OBL-9 · CI-1 · $4.00", MUTED),
    ]
    for i, (cb, label, col) in enumerate(items):
        card_rect(s, 0.55, 2.35 + i*0.45, 5.5, 0.38,
                  fill=RGBColor(0x1a, 0x24, 0x38) if cb == "☑" else SURF2,
                  border=BLUE if cb == "☑" else BORDER)
        tb(s, f"{cb}  {label}", l=0.7, t=2.4 + i*0.45, w=5.2, h=0.28,
           size=9.5, color=col)

    form_fields = [
        ("Refund amount:", "$40.00"),
        ("Method:",        "Stripe (original payment method)"),
        ("Reason:",        "Duplicate charge"),
    ]
    for i, (label, val) in enumerate(form_fields):
        y = 3.35 + i * 0.42
        tb(s, label, l=0.55, t=y, w=1.7, h=0.3, size=9.5, color=MUTED)
        tb(s, val,   l=2.3,  t=y, w=3.8, h=0.3, size=9.5, color=TEXT)

    card_rect(s, 2.5, 4.75, 2.5, 0.42, fill=BLUE, border=BLUE)
    tb(s, "Submit Refund →", l=2.5, t=4.8, w=2.5, h=0.32,
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # What happens steps
    tb(s, "What happens after Submit:", l=6.4, t=1.25, w=6.7, h=0.28,
       size=11, bold=True, color=TEXT)

    steps = [
        ("①", "Refund record created → 202",
         "UI shows 'Refund Requested'. AuditLog entry: refund_created.", BLUE),
        ("②", "Refund Processor dequeues",
         "Calls VendorAdapter.refund(cart, items) → Stripe API.", MUTED),
        ("③", "Stripe emits charge.refunded",
         "EventProcessor receives as payment_refunded VendorEvent.", GOLD),
        ("④", "Ledger appends LDG-7 (refunded)",
         "Allocation: -$40 against OBL-1. outstanding_amount → $40 → status OPEN.", GREEN),
        ("⑤", "Cart.refund_amount += $40",
         "Cart stays payment_settled. Refund is a ledger event, not a status change.", TEXT),
    ]

    for i, (num, title, desc, col) in enumerate(steps):
        y = 1.62 + i * 0.95
        card_rect(s, 6.4, y, 6.7, 0.85, fill=SURF)
        tb(s, num, l=6.55, t=y+0.1, w=0.3, h=0.3, size=13, bold=True, color=col)
        tb(s, title, l=6.9, t=y+0.07, w=6.0, h=0.3, size=10, bold=True, color=col)
        tb(s, desc,  l=6.9, t=y+0.38, w=6.0, h=0.35, size=9, color=MUTED)

    card_rect(s, 6.4, 6.52, 6.7, 0.42, fill=SURF2)
    tb(s, "Cart items are refunded in full only. Pick specific items to achieve a partial refund of the overall cart.",
       l=6.55, t=6.58, w=6.5, h=0.35, size=9, color=MUTED)


def slide_21_obligation_adjustment(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Support UI — Obligation Adjustment")

    # Adjustment form
    card_rect(s, 0.4, 1.2, 5.8, 5.4, fill=SURF)
    tb(s, "Adjustment Form (OBL-5 · 000328981 · citation_fee · $40.00)",
       l=0.55, t=1.28, w=5.5, h=0.28, size=10, bold=True, color=TEXT)

    tb(s, "Current State:", l=0.55, t=1.62, w=5.5, h=0.25, size=9.5, bold=True, color=MUTED)
    current_fields = [
        ("id:", "5 (000328981)"),
        ("label:", "citation_fee"),
        ("amount:", "$40.00"),
        ("allocated_total:", "$0.00"),
        ("outstanding_amount:", "$40.00"),
        ("status:", "OPEN"),
    ]
    for i, (label, val) in enumerate(current_fields):
        y = 1.92 + i * 0.32
        tb(s, label, l=0.7, t=y, w=1.6, h=0.27, size=9, color=MUTED)
        tb(s, val,   l=2.4, t=y, w=3.7, h=0.27, size=9, color=TEXT, font="Courier New")

    tb(s, "Select Action:", l=0.55, t=4.0, w=5.5, h=0.25, size=9.5, bold=True, color=TEXT)
    actions = [
        ("☑", "Waive",        "Forgive full $40.00 (unpaid, no prior allocation)", True),
        ("☐", "Partial Waive","Requires prior partial payment (disabled)", False),
        ("☐", "Supersede",    "Replace with new amount", True),
    ]
    for i, (cb, name, desc, enabled) in enumerate(actions):
        y = 4.3 + i * 0.42
        col = GOLD if (cb == "☑" and enabled) else (MUTED if not enabled else TEXT)
        tb(s, f"{cb}  {name} — {desc}", l=0.7, t=y, w=5.2, h=0.3, size=9, color=col)

    tb(s, "Reason: Court-ordered hardship waiver", l=0.7, t=5.65, w=5.0, h=0.28, size=9, color=MUTED)

    card_rect(s, 2.2, 6.1, 2.5, 0.35, fill=RGBColor(0x0a, 0x14, 0x0a), border=GREEN)
    tb(s, "Apply Waiver →", l=2.2, t=6.14, w=2.5, h=0.27, size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # State matrix
    tb(s, "Available actions by state:", l=6.5, t=1.2, w=6.5, h=0.28, size=11, bold=True, color=TEXT)
    state_headers = ["State", "Waive", "Partial Waive", "Supersede"]
    state_rows = [
        ["OPEN (unpaid)",    "✓ Available", "✕ Not available", "✓ Available"],
        ["OPEN (partial paid)", "✕",       "✓ Available",      "✕"],
        ["LOCKED",           "✕",          "✕",                "✕"],
        ["CLOSED",           "✕",          "✕",                "✕"],
        ["WAIVED",           "✕",          "✕",                "✕"],
    ]
    add_table(s, state_headers, state_rows, l=6.5, t=1.55, w=6.5, h=2.0,
              font_size=9, col_widths=[0.28, 0.24, 0.24, 0.24])

    # Code block
    card_rect(s, 6.5, 3.7, 6.5, 3.0, fill=SURF2)
    tb(s, "What each action writes:", l=6.65, t=3.75, w=6.2, h=0.25, size=9.5, bold=True, color=TEXT)
    code = """# Waive (unpaid)
obligation.waived_amount  += $40.00
obligation.outstanding    =  $0.00
obligation.status         =  WAIVED (terminal)

# Supersede (creates new obligation)
old_obligation.status     =  SUPERSEDED (terminal)
new_obligation.amount     =  new_amount
new_obligation.status     =  OPEN"""
    tb(s, code, l=6.65, t=4.05, w=6.2, h=2.5, size=9, color=TEXT, font="Courier New")


def slide_22_full_data_transition(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Online Payment — Full Data Transition (Stripe)",
                subtitle="Each step shows the exact state of every entity. All transitions are driven by events — the API returns immediately at step 3.")

    steps = [
        ("① POST /checkout\n11:56:23Z", BLUE,
         "CART 1: draft · $84.00\nCI-2→OBL-1 $40\nCI-3→OBL-2 $40\nCI-1→OBL-9 $4\n\nOBL-1: OPEN · $40 owed\nOBL-2: OPEN · $40 owed\n\nLEDGER: empty"),
        ("② Checkout confirmed", MUTED,
         "CART 1: checkout\nvendor_intent_id:\npi_3RKstt…\n\nOBL-1: OPEN\nOBL-2: OPEN\n\nLEDGER: empty"),
        ("③ POST /submit →202\n11:56:23Z", GOLD,
         "CART 1: payment_submitted\n\nOBL-1: LOCKED\n  locked_by: cart:1\nOBL-2: LOCKED\n  locked_by: cart:1\n\nLEDGER: empty\n\n← API returns here\nUser sees 'Processing…'"),
        ("④ payment_confirmed\n11:56:29Z", BLUE,
         "CART 1: payment_confirmed\n\nOBL-1: LOCKED\n  allocated: $40\n  outstanding: $0\nOBL-2: LOCKED · $40\n\nLEDGER (event:1):\nLDG-1 CONFIRMED $84\n  →OBL-1 $40\n  →OBL-2 $40\n  →OBL-9 $4\nLDG-2 fee -$1.55"),
        ("⑤ payment_settled\n11:56:33Z", GREEN,
         "CART 1: settled\nrefund_amount: $0\n\nOBL-1: CLOSED\n  allocated: $40\nOBL-2: CLOSED · $40\n\nLEDGER (event:2):\nLDG-1 SETTLED\n  txn_3RKstt…\nLDG-2 SETTLED\n\n← In financial reports"),
    ]

    for i, (title, col, content) in enumerate(steps):
        x = 0.35 + i * 2.6
        card_rect(s, x, 1.1, 2.45, 6.1, fill=SURF)
        tb(s, title, l=x+0.1, t=1.18, w=2.25, h=0.55,
           size=9.5, bold=True, color=col)
        tb(s, content, l=x+0.1, t=1.78, w=2.25, h=5.3,
           size=8.5, color=TEXT, font="Courier New")


def slide_23_cart_activity(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Support UI — Cart Activity Log",
                subtitle="Append-only log of every cart status transition. Linked to obligations and vendor events via correlation_id.")

    tb(s, "Cart Activity Log (Cart 1):", l=0.4, t=1.1, w=8, h=0.28, size=11, bold=True, color=TEXT)

    cal_headers = ["timestamp", "event", "from", "to", "actor", "cause"]
    cal_rows = [
        ["11:56:15Z", "cart.created",            "—",         "draft",     "demo_user", "req_checkout_..."],
        ["11:56:15Z", "cart.items_added",         "draft",     "draft",     "demo_user", "req_checkout_..."],
        ["11:56:23Z", "cart.checked_out",         "draft",     "checkout",  "demo_user", "req_checkout_..."],
        ["11:56:23Z", "cart.payment_submitted",   "checkout",  "submitted", "demo_user", "req_submit_1_..."],
        ["11:56:29Z", "cart.payment_confirmed",   "submitted", "confirmed", "worker",    "vendor_event:1"],
        ["11:56:33Z", "cart.payment_settled",     "confirmed", "settled",   "worker",    "vendor_event:2"],
    ]
    add_table(s, cal_headers, cal_rows, l=0.4, t=1.45, w=12.33, h=2.2,
              font_size=9.5, col_widths=[0.12, 0.22, 0.1, 0.12, 0.13, 0.31])

    # Record structure
    tb(s, "CartActivityLog Record Structure:", l=0.4, t=3.85, w=6, h=0.28,
       size=10, bold=True, color=TEXT)
    card_rect(s, 0.4, 4.18, 6.0, 3.1, fill=SURF2)
    record_code = """cart_id:           1
event_type:        cart.payment_settled
old_status:        payment_confirmed
new_status:        payment_settled
actor_type:        worker
actor_id:          worker_event:2
correlation_id:    cart_payment:1
causation_id:      vendor_event:2
metadata:          {"vendor": "stripe", "amount": 8400}
timestamp:         2026-03-24T11:56:33Z"""
    tb(s, record_code, l=0.55, t=4.28, w=5.7, h=2.9,
       size=9, color=TEXT, font="Courier New")

    # Support use
    card_rect(s, 6.6, 3.85, 6.13, 3.43, fill=RGBColor(0x0d, 0x18, 0x0d), border=GREEN)
    tb(s, "Support Use", l=6.75, t=3.92, w=5.8, h=0.28, size=10, bold=True, color=GREEN)
    tb(s, "If a violator says 'I paid but my citation still shows unpaid', "
          "filter this log by citation → find cart → check whether it reached payment_settled.\n\n"
          "If stuck at payment_submitted, the vendor event has not arrived yet.\n\n"
          "Every cart transition is linked back to the exact request or vendor event that caused it — "
          "support staff can trace any payment complaint in one query.",
       l=6.75, t=4.25, w=5.8, h=2.9, size=9.5, color=TEXT)


def slide_24_obligation_activity(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Support UI — Obligation Activity Log",
                subtitle="Every balance change and status transition on each obligation — linked to the cart and vendor event that caused it.")

    tb(s, "Obligation Activity Log (OBL-1 · 000331714 · citation_fee):",
       l=0.4, t=1.1, w=12.33, h=0.28, size=11, bold=True, color=TEXT)

    oal_headers = ["timestamp", "event", "status change", "outstanding", "allocated", "cause"]
    oal_rows = [
        ["11:56:15Z", "obligation.created",            "— → OPEN",           "$40.00", "$0.00",  "citation_issued"],
        ["11:56:23Z", "obligation.locked",             "OPEN → LOCKED",       "$40.00", "$0.00",  "cart_item:2"],
        ["11:56:29Z", "obligation.payment_confirmed",  "LOCKED → CLOSED",     "$0.00",  "$40.00", "vendor_event:1"],
        ["11:56:33Z", "obligation.payment_settled",    "CLOSED → CLOSED",     "$0.00",  "$40.00", "vendor_event:2"],
        ["11:57:02Z", "obligation.payment_settled",    "CLOSED → CLOSED",     "$0.00",  "$80.00 +overpaid", "checkalt 41/1"],
        ["11:57:21Z", "obligation.payment_refunded",   "CLOSED → CLOSED",     "$0.00",  "$40.00 (−$40)", "vendor_event:5"],
    ]
    add_table(s, oal_headers, oal_rows, l=0.4, t=1.45, w=12.33, h=2.2,
              font_size=9, col_widths=[0.12, 0.24, 0.18, 0.14, 0.18, 0.14])

    # Record structure
    tb(s, "ObligationActivityLog Record Structure:", l=0.4, t=3.85, w=6, h=0.28,
       size=10, bold=True, color=TEXT)
    card_rect(s, 0.4, 4.18, 6.0, 3.1, fill=SURF2)
    code = """obligation_id:     1 (000331714)
event_type:        obligation.payment_refunded
old_status:        CLOSED / new_status: CLOSED
old_outstanding:   0  / new_outstanding:  0
old_allocated:     8000 / new_allocated: 4000
allocated_delta:   -4000
ledger_entry_id:   LDG-7
correlation_id:    cart_payment:1
causation_id:      vendor_event:5
timestamp:         2026-03-24T11:57:21Z"""
    tb(s, code, l=0.55, t=4.28, w=5.7, h=2.9, size=9, color=TEXT, font="Courier New")

    card_rect(s, 6.6, 3.85, 6.13, 3.43, fill=RGBColor(0x0d, 0x18, 0x0d), border=GREEN)
    tb(s, "Support Use", l=6.75, t=3.92, w=5.8, h=0.28, size=10, bold=True, color=GREEN)
    tb(s, "Answer 'why is this obligation still open?' in one query.\n\n"
          "The activity log shows every allocation added and reversed, "
          "the ledger entry that caused it, and who/what initiated the change.\n\n"
          "Each record has an allocated_delta showing exactly how much was added "
          "or removed in that event — no manual arithmetic required.",
       l=6.75, t=4.25, w=5.8, h=2.9, size=9.5, color=TEXT)


def slide_25_ledger_report(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Support UI — Ledger Report",
                subtitle="Each ledger item covers a cart-level charge. Allocations within each item map the gross amount to obligations (1 item → n allocations).")

    # Filters bar
    card_rect(s, 0.4, 1.08, 12.33, 0.38, fill=SURF2)
    tb(s, "Date: 2026-03-24    Vendor: All    Mode: All    [Export CSV]",
       l=0.55, t=1.13, w=12.0, h=0.28, size=9, color=MUTED)

    # Table
    headers = ["ldg_id", "type", "label", "cart", "vendor", "amount", "stage", "vendor_ref"]
    rows = [
        ["LDG-1",   "payment", "cart_payment",     "1", "stripe · online", "$84.00",  "settled",  "txn_3RKstt…S1G9OXSH1"],
        ["  ↳ CI-2","alloc",   "citation_fee",      "",  "OBL-1 · 000331714", "$40.00", "",        ""],
        ["  ↳ CI-3","alloc",   "citation_fee",      "",  "OBL-2 · 000329181", "$40.00", "",        ""],
        ["  ↳ CI-1","alloc",   "convenience_fee",   "",  "OBL-9 · cart:1",   "$4.00",  "",        ""],
        ["LDG-2",   "fee",     "stripe_fee",        "1", "stripe",           "-$1.55", "settled",  "txn_3RKstt…S1G9OXSH1"],
        ["  (no allocations — fee charged at cart level)", "", "", "", "", "", "", ""],
        ["LDG-7",   "payment", "refund_payment",    "1", "stripe · online",  "-$40.00","refunded", "txn_3RKstt…S138WeBtC"],
        ["  ↳ CI-2","alloc",   "citation_fee",      "",  "OBL-1 · 000331714","-$40.00","",        ""],
        ["LDG-6",   "payment", "check_payment",     "4", "checkalt · check", "$20.00", "settled",  "42/1"],
        ["  ↳ CI-9","alloc",   "citation_fee",      "",  "OBL-3 · 000322452","$20.00", "",        ""],
    ]
    add_table(s, headers, rows, l=0.4, t=1.52, w=12.33, h=4.0,
              font_size=8.5, col_widths=[0.10, 0.09, 0.15, 0.06, 0.22, 0.09, 0.10, 0.19])

    # Summary boxes
    summaries = [
        ("Payments (Gross)", "$104.00", "LDG-1 $84 + LDG-6 $20", TEXT),
        ("Refunds",          "-$40.00", "LDG-7",                  RED),
        ("Vendor Fees",      "-$1.55",  "LDG-2",                  MUTED),
        ("Net Settled",      "$62.45",  "$104 − $40 − $1.55",     GREEN),
    ]
    for i, (label, amount, note, col) in enumerate(summaries):
        x = 0.4 + i * 3.1
        card_rect(s, x, 5.7, 2.9, 1.0, fill=SURF, border=col)
        tb(s, label,  l=x+0.12, t=5.77, w=2.66, h=0.28, size=8.5, color=MUTED)
        tb(s, amount, l=x+0.12, t=6.07, w=2.66, h=0.35, size=16, bold=True, color=col)
        tb(s, note,   l=x+0.12, t=6.45, w=2.66, h=0.2,  size=8, color=MUTED)


def slide_26_triage(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Support UI — Triage System",
                subtitle="Centralized queue for support staff to flag, investigate, and resolve payment anomalies.")

    tb(s, "Triage Queue (4 open items):", l=0.4, t=1.08, w=12.33, h=0.28,
       size=11, bold=True, color=TEXT)

    triage_headers = ["priority", "type", "citation", "cart", "detail", "assigned", "status"]
    triage_rows = [
        ["HIGH", "dispute",         "000330575", "2", "Chargeback filed $84.00",             "sarah.k",  "in_progress"],
        ["MED",  "overpayment",     "000329181", "1,3","Paid by Stripe + CheckAlt ($40 over)","—",       "open"],
        ["MED",  "partial_payment", "000322452", "4", "Check paid $20, still owes $20",      "james.r",  "in_progress"],
        ["LOW",  "refund_open",     "000331714", "1", "Stripe refund issued, OBL-1 still OPEN","—",      "open"],
    ]
    add_table(s, triage_headers, triage_rows, l=0.4, t=1.42, w=12.33, h=1.5,
              font_size=9, col_widths=[0.09, 0.15, 0.13, 0.07, 0.31, 0.13, 0.12])

    # Item detail
    card_rect(s, 0.4, 3.1, 5.9, 1.85, fill=SURF2)
    tb(s, "Triage Item Detail: CRT-449 — Stuck Payment", l=0.55, t=3.15, w=5.6, h=0.28,
       size=10, bold=True, color=GOLD)
    detail_code = """Cart submitted:    Mar 13 09:22Z
Status:            payment_submitted (48+ hrs)
Stripe intent:     pi_3Px7Rb... → requires_action
Last access log:   req_4mNx1k (submit)
Suggested action:  Cancel intent + reopen obligations"""
    tb(s, detail_code, l=0.55, t=3.47, w=5.6, h=1.4,
       size=8.5, color=TEXT, font="Courier New")

    # Auto-detection table
    tb(s, "Triage item types & auto-detection:", l=0.4, t=5.1, w=12.33, h=0.28,
       size=10, bold=True, color=TEXT)
    detect_headers = ["Type", "Detection Rule", "Action"]
    detect_rows = [
        ["dispute",              "VendorEvent payment_dispute_funds_withdrawn received",     "Auto-create HIGH priority. Track evidence window."],
        ["stuck_payment",        "Cart in payment_submitted > 24 hrs",                       "Check Stripe intent status. Cancel + reopen if needed."],
        ["refund_failed",        "Refund processor errors ≥ 3 attempts",                    "Manual Stripe refund. Mark processed."],
        ["balance_mismatch",     "Obligation outstanding > 0 but cart settled",              "Inspect ledger entries. Replay or adjust manually."],
        ["vendor_event_failed",  "VendorEvent processing_status = failed",                  "Inspect error, fix root cause, replay event."],
    ]
    add_table(s, detect_headers, detect_rows, l=0.4, t=5.43, w=12.33, h=2.0,
              font_size=9, col_widths=[0.17, 0.4, 0.43])


def slide_27_comparison(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Current vs Proposed — Head-to-Head")

    headers = ["Capability", "Current System (Victor + IMP)", "Gringotts"]
    rows = [
        ["Financial source of truth",
         "Mutable Citations fields — can drift from actual payments",
         "Immutable ledger — every balance derived from append-only entries"],
        ["Refund handling",
         "citation_amount_paid never decremented. Balance wrong after refund.",
         "Refund event creates negative ledger allocation. Obligation reopens automatically."],
        ["Check amounts",
         "Always $40 (hardcoded in IMP regardless of actual check)",
         "Actual amount from vendor event data (no hardcoding)"],
        ["Partial payments",
         "Not supported for checks. Stripe tracks PARTIALLY_PAID but no history.",
         "Native support. Multiple payments accumulate in ledger allocations."],
        ["Dispute (chargeback)",
         "No chargeback lifecycle. Refund tracked but balance not reversed.",
         "Full disputed → won/lost lifecycle with ledger stage tracking."],
        ["Audit trail",
         "IMP creates activity logs. Victor does not. S3 has raw JSONs.",
         "Every state change logged with before/after + correlation/causation IDs."],
        ["Idempotency",
         "Victor dedupes on stripe_event_id. IMP has no deduplication.",
         "All event processors deduplicate. Ledger identifies duplicate entries."],
        ["Multi-citation checkout",
         "One Stripe charge per citation",
         "One cart = one Stripe charge for N citations. Lower fees."],
        ["In-person payment record",
         "No financial record created. Invisible to support and reporting.",
         "VendorEvent + ledger entry created. Fully auditable."],
        ["Financial reporting",
         "Report Manager joins multiple tables that can drift. In-person never appears.",
         "Single ledger query. All payment modes in one consistent source."],
    ]
    add_table(s, headers, rows, l=0.4, t=1.1, w=12.33, h=6.1,
              font_size=9, col_widths=[0.22, 0.39, 0.39])


def slide_28_migration(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Migration Strategy")

    principles = [
        "No big-bang cutover — phased by payment mode",
        "Run in parallel during transition (dual-write if needed)",
        "Existing citations backfilled into Gringotts obligations",
        "Victor and IMP remain active for legacy citations until fully migrated",
    ]
    card_rect(s, 0.4, 1.1, 12.33, 0.85, fill=SURF)
    for i, p in enumerate(principles):
        tb(s, "·  " + p, l=0.6, t=1.15 + i*0.2, w=12.0, h=0.2,
           size=9.5, color=MUTED)

    phases = [
        ("Phase 1", "Gringotts Infrastructure", BLUE, [
            "Deploy Gringotts alongside existing services — Ledger, Obligation, Cart, VendorEvent tables created. No traffic yet.",
            "Backfill existing citations as Obligations — Open citations → OPEN obligations. Paid citations → CLOSED obligations with synthetic ledger entries.",
        ]),
        ("Phase 2", "Online Payments (Stripe)", GOLD, [
            "Route new Stripe webhooks to Gringotts — Victor reads events; Gringotts processes them. Dual-write until validated.",
            "New checkouts create Gringotts Carts — Public API updated to use POST /checkout + POST /submit.",
            "Decommission Victor's payment logic — Victor becomes a thin webhook receiver only.",
        ]),
        ("Phase 3", "Check Payments + Reports", GREEN, [
            "Route CheckAlt events to Gringotts — IMP replaced by Gringotts CheckAlt vendor module.",
            "Replace Report Manager with ledger queries — GET /reports/ledger replaces CSV generation.",
            "Decommission IMP + Report Manager.",
        ]),
    ]

    for i, (phase, title, col, steps) in enumerate(phases):
        x = 0.4 + i * 4.28
        card_rect(s, x, 2.1, 4.0, 5.1, fill=SURF, border=col)
        badge = s.shapes.add_shape(1, Inches(x+0.15), Inches(2.18),
                                   Inches(1.1), Inches(0.28))
        badge.fill.solid(); badge.fill.fore_color.rgb = col
        badge.line.fill.background()
        tb(s, phase, l=x+0.17, t=2.18, w=1.06, h=0.28,
           size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)
        tb(s, title, l=x+0.15, t=2.52, w=3.7, h=0.3, size=11, bold=True, color=col)
        step_lines = []
        for j, step in enumerate(steps):
            step_lines.append({"text": f"{j+1}.  {step}", "size": 9.5, "color": TEXT})
            step_lines.append({"text": "", "size": 4, "color": TEXT})
        mtb(s, step_lines, l=x+0.15, t=2.88, w=3.7, h=4.0)


def slide_29_summary(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Summary & Next Steps")

    # What Gringotts Fixes
    tb(s, "What Gringotts Fixes", l=0.4, t=1.1, w=5.9, h=0.3,
       size=13, bold=True, color=GOLD)
    card_rect(s, 0.4, 1.48, 5.9, 4.3, fill=SURF)
    fixes = [
        "Immutable ledger replaces mutable Citations fields\n  — Financial history is complete, reconstructable, and consistent.",
        "Real check amounts from CheckAlt\n  — No more hardcoded $40. Obligations reflect what was actually paid.",
        "Refunds, disputes, waivers fully modeled\n  — Every financial operation has a defined, auditable lifecycle.",
        "One Stripe charge for multiple citations\n  — Cart groups N obligations into one charge. Lower fees.",
        "Support UI with full traceability\n  — Waive, supersede, or refund directly — with full audit trail.",
    ]
    fix_lines = []
    for f in fixes:
        fix_lines.append({"text": "✓  " + f.split("\n")[0], "size": 10, "bold": True, "color": GREEN})
        fix_lines.append({"text": "     " + f.split("\n")[1].strip(), "size": 9.5, "color": MUTED})
        fix_lines.append({"text": "", "size": 4, "color": TEXT})
    mtb(s, fix_lines, l=0.55, t=1.58, w=5.6, h=4.1)

    # Next Steps
    tb(s, "Next Steps", l=6.5, t=1.1, w=6.5, h=0.3,
       size=13, bold=True, color=BLUE)
    card_rect(s, 6.5, 1.48, 6.5, 2.8, fill=SURF)
    next_steps = [
        "1.  Finalize Obligation backfill script for existing citations",
        "2.  Define dual-write strategy for Stripe webhooks (Victor ↔ Gringotts)",
        "3.  Build Support UI (Phase 1: citation lookup + ledger view)",
        "4.  Validate check amounts against CheckAlt source data",
        "5.  Reconcile Report Manager output vs Gringotts ledger queries",
        "6.  Set go-live criteria per payment mode",
    ]
    ns_lines = [{"text": step, "size": 10, "color": TEXT} for step in next_steps]
    mtb(s, ns_lines, l=6.65, t=1.58, w=6.2, h=2.6)

    # Open Questions
    tb(s, "Open Questions for Discussion", l=6.5, t=4.45, w=6.5, h=0.3,
       size=12, bold=True, color=PURPLE)
    card_rect(s, 6.5, 4.82, 6.5, 2.5, fill=RGBColor(0x14, 0x0d, 0x1f), border=PURPLE)
    questions = [
        "What is the target timeline for Phase 2 cutover?",
        "Do we need dual-write or can we trust backfill + cutover?",
        "Who owns the Support UI build — backend or frontend team?",
        "Should Report Manager be deprecated or kept for legacy reports?",
    ]
    q_lines = [{"text": "?  " + q, "size": 10, "color": PURPLE} for q in questions]
    mtb(s, q_lines, l=6.65, t=4.92, w=6.2, h=2.3)


# ═══════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    print("Building slides...")
    slide_01_title(prs)             ; print("  1/29 Title")
    slide_02_agenda(prs)            ; print("  2/29 Agenda")
    slide_03_divider(prs, "01", "Current System",
                     "A review of how citation payments work today across Victor, IMP, and Report Manager\n— and the hidden risks buried in each service.")
    print("  3/29 Section: Current System")
    slide_04_services(prs)          ; print("  4/29 Services & Payment Paths")
    slide_05_how_payments(prs)      ; print("  5/29 How Payments Work Today")
    slide_06_whats_broken(prs)      ; print("  6/29 What's Broken")
    slide_07_gringotts_divider(prs) ; print("  7/29 Section: Gringotts")
    slide_08_philosophy(prs)        ; print("  8/29 Design Philosophy")
    slide_09_data_models(prs)       ; print("  9/29 Core Data Models")
    slide_10_stripe_flow(prs)       ; print(" 10/29 Online Payment Flow (Stripe)")
    slide_11_check_flow(prs)        ; print(" 11/29 Check & In-Person Flow")
    slide_12_ledger_truth(prs)      ; print(" 12/29 Ledger as Source of Truth")
    slide_13_obligation_ops(prs)    ; print(" 13/29 Obligation Operations")
    slide_14_architecture(prs)      ; print(" 14/29 System Architecture")
    slide_15_ui_divider(prs)        ; print(" 15/29 Section: Support UI")
    slide_16_citation_overview(prs) ; print(" 16/29 Citation Payment Overview")
    slide_17_sample_data(prs)       ; print(" 17/29 Obligation & Cart Sample Data")
    slide_18_ledger_transitions(prs); print(" 18/29 Ledger Data Transitions")
    slide_19_access_audit(prs)      ; print(" 19/29 Access Log & Audit Log")
    slide_20_refund_ui(prs)         ; print(" 20/29 Initiating a Refund")
    slide_21_obligation_adjustment(prs); print(" 21/29 Obligation Adjustment")
    slide_22_full_data_transition(prs) ; print(" 22/29 Full Data Transition (Stripe)")
    slide_23_cart_activity(prs)     ; print(" 23/29 Cart Activity Log")
    slide_24_obligation_activity(prs); print(" 24/29 Obligation Activity Log")
    slide_25_ledger_report(prs)     ; print(" 25/29 Ledger Report")
    slide_26_triage(prs)            ; print(" 26/29 Triage System")
    slide_27_comparison(prs)        ; print(" 27/29 Current vs Proposed")
    slide_28_migration(prs)         ; print(" 28/29 Migration Strategy")
    slide_29_summary(prs)           ; print(" 29/29 Summary & Next Steps")

    out = "/Users/obvio/Documents/code_repos/gringotts/gringotts_presentation.pptx"
    prs.save(out)
    print(f"\nSaved → {out}")


if __name__ == "__main__":
    main()
