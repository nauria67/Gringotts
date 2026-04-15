"""
Gringotts Support Team Guide — PowerPoint Generator
Requires: pip install python-pptx
Run:      python3 make_support_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Colors ────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0d, 0x11, 0x17)
SURF    = RGBColor(0x16, 0x1b, 0x22)
SURF2   = RGBColor(0x21, 0x26, 0x2d)
BORDER  = RGBColor(0x30, 0x36, 0x3d)
TEXT    = RGBColor(0xe6, 0xed, 0xf3)
MUTED   = RGBColor(0x8b, 0x94, 0x9e)
GOLD    = RGBColor(0xf0, 0xa5, 0x00)
BLUE    = RGBColor(0x58, 0xa6, 0xff)
GREEN   = RGBColor(0x3f, 0xb9, 0x50)
RED     = RGBColor(0xff, 0x7b, 0x72)
PURPLE  = RGBColor(0xd2, 0xa8, 0xff)
WHITE   = RGBColor(0xff, 0xff, 0xff)

BADGE_CLOSED   = RGBColor(0x0f, 0x2a, 0x18)
BADGE_OPEN     = RGBColor(0x0d, 0x1f, 0x38)
BADGE_LOCKED   = RGBColor(0x2e, 0x20, 0x00)
BADGE_WAIVED   = RGBColor(0x22, 0x14, 0x38)
BADGE_DISPUTED = RGBColor(0x38, 0x0e, 0x0e)

W = Inches(13.33)
H = Inches(7.5)


# ─── Helpers ───────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, text, l, t, w, h,
       size=14, bold=False, color=TEXT,
       align=PP_ALIGN.LEFT, italic=False,
       font="Calibri", wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font
    return txb


def mtb(slide, lines, l, t, w, h,
        default_size=12, default_color=TEXT,
        default_font="Calibri", wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = line.get("align", PP_ALIGN.LEFT)
        if line.get("space_before"):
            p.space_before = Pt(line["space_before"])
        run = p.add_run()
        run.text         = line.get("text", "")
        run.font.size    = Pt(line.get("size",   default_size))
        run.font.bold    = line.get("bold",   False)
        run.font.italic  = line.get("italic", False)
        run.font.color.rgb = line.get("color", default_color)
        run.font.name    = line.get("font",   default_font)
    return txb


def slide_title(slide, text, subtitle=None):
    rect = slide.shapes.add_shape(
        1, Inches(0.5), Inches(0.35), Inches(12.33), Inches(0.04))
    rect.fill.solid()
    rect.fill.fore_color.rgb = GOLD
    rect.line.fill.background()
    tb(slide, text, l=0.5, t=0.42, w=12.33, h=0.55,
       size=22, bold=True, color=GOLD)
    if subtitle:
        tb(slide, subtitle, l=0.5, t=0.92, w=12.33, h=0.35,
           size=11, color=MUTED)


def card_rect(slide, l, t, w, h, fill=SURF, border=BORDER):
    shape = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    return shape


def add_table(slide, headers, rows, l, t, w, h,
              font_size=9, col_widths=None):
    ncols = len(headers)
    nrows = len(rows) + 1
    table = slide.shapes.add_table(
        nrows, ncols,
        Inches(l), Inches(t), Inches(w), Inches(h)).table

    if col_widths:
        total = Inches(w)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = int(total * cw)

    STATUS_COLORS = {
        "CLOSED": GREEN, "OPEN": BLUE, "LOCKED": GOLD,
        "WAIVED": PURPLE, "SUPERSEDED": MUTED, "VOIDED": MUTED,
        "DISPUTED": RED, "dispute_lost": RED,
        "payment_settled": GREEN, "payment_confirmed": BLUE,
        "payment_submitted": GOLD, "checkout": GOLD,
        "draft": MUTED, "disputed": RED, "abandoned": MUTED,
        "open": BLUE, "in_progress": GOLD, "resolved": GREEN,
        "HIGH": RED, "MED": GOLD, "LOW": MUTED,
    }

    def style_cell(cell, text, bold=False, bg=SURF, size=None):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size or font_size)
        run.font.bold  = bold
        run.font.name  = "Calibri"
        col = MUTED if bold else STATUS_COLORS.get(text, TEXT)
        if not bold:
            if text.startswith("✓"): col = GREEN
            elif text.startswith("✕"): col = RED
            elif text.startswith("HIGH"): col = RED
            elif text.startswith("MED"):  col = GOLD
            elif text.startswith("LOW"):  col = MUTED
        run.font.color.rgb = col
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

    for ci, hdr in enumerate(headers):
        style_cell(table.cell(0, ci), hdr, bold=True, bg=SURF2)
    for ri, row in enumerate(rows):
        bg = SURF if ri % 2 == 0 else SURF2
        for ci, val in enumerate(row):
            style_cell(table.cell(ri + 1, ci), str(val), bg=bg)

    return table


def section_divider(prs, num, title, desc, accent=GOLD):
    s = blank_slide(prs)
    set_bg(s)
    tb(s, num, l=8.5, t=0.5, w=4.5, h=4.0,
       size=200, bold=True, color=RGBColor(0x1a, 0x20, 0x28),
       align=PP_ALIGN.RIGHT)
    bar = s.shapes.add_shape(
        1, Inches(0.5), Inches(2.0), Inches(0.08), Inches(2.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    tb(s, title, l=0.75, t=2.1, w=9, h=1.6,
       size=52, bold=True, color=TEXT)
    tb(s, desc, l=0.75, t=3.9, w=8.5, h=1.0,
       size=14, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    s = blank_slide(prs)
    set_bg(s)
    bar = s.shapes.add_shape(
        1, Inches(0), Inches(3.05), Inches(0.12), Inches(1.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    tb(s, "Support Team Guide  ·  April 2026",
       l=0.6, t=1.5, w=10, h=0.4, size=12, color=MUTED)
    tb(s, "Gringotts",
       l=0.6, t=1.95, w=10, h=1.1, size=64, bold=True, color=GOLD)
    tb(s, "Payment Lookup  ·  Refunds  ·  Obligation Adjustments  ·  Triage",
       l=0.6, t=3.15, w=11, h=0.45, size=16, color=TEXT)
    tb(s, "This guide covers everything you need to look up citation payments, "
          "investigate issues, take action on obligations, and work the triage queue. "
          "No engineering background required.",
       l=0.6, t=3.75, w=9.5, h=0.85, size=11, color=MUTED)


def slide_02_agenda(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "What's in This Guide")

    sections = [
        ("01", "How It Works",
         "What Gringotts does, how a payment travels from citation to settled, and the three key concepts.",
         BLUE),
        ("02", "Reading a Citation",
         "How to look up a citation, interpret obligation and cart statuses, and understand the numbers.",
         GREEN),
        ("03", "Actions You Can Take",
         "When and how to issue a refund, waive an obligation, or supersede an obligation.",
         GOLD),
        ("04", "Investigating Complaints",
         "How to trace a payment issue from complaint to root cause using audit logs.",
         PURPLE),
        ("05", "The Triage Queue",
         "Your daily work queue: item types, priorities, auto-detection, and how to resolve each.",
         RED),
    ]
    for i, (num, title, desc, color) in enumerate(sections):
        x = 0.4 + i * 2.5
        card_rect(s, x, 1.35, 2.35, 5.3, fill=SURF, border=color)
        tb(s, num, l=x+0.12, t=1.48, w=2.1, h=0.55,
           size=30, bold=True, color=color)
        tb(s, title, l=x+0.12, t=2.1, w=2.11, h=0.45,
           size=12, bold=True, color=TEXT)
        tb(s, desc, l=x+0.12, t=2.6, w=2.11, h=3.8,
           size=9, color=MUTED)


# ─── SECTION 01 ────────────────────────────────────────────────────────────

def slide_04_what_is_gringotts(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "What is Gringotts?",
                subtitle="Gringotts is the system that records, tracks, and manages all citation payments.")

    panels = [
        ("Citations Create Obligations",
         "When a traffic citation is issued, it creates an obligation — a record of what is owed. "
         "Gringotts tracks how much has been paid, what's still outstanding, and "
         "whether any amounts were waived or adjusted by support.",
         BLUE,
         ["citation_fee    $40.00", "status          OPEN", "outstanding     $40.00"]),
        ("Every Payment Is Recorded Permanently",
         "Every payment — online card, check, or in-person — is recorded in a ledger "
         "that can never be modified. All outstanding balances are calculated from "
         "this ledger. Nothing is ever erased or overwritten.",
         GREEN,
         ["LDG-1  +$40.00  payment_settled", "  → OBL-1  citation_fee  $40", "LDG-7  -$40.00  payment_refunded"]),
        ("You Handle What the System Can't",
         "Gringotts processes normal payments automatically. "
         "When something unusual happens — a chargeback, an overpayment, a stuck payment — "
         "it appears in the Triage Queue for you to investigate and resolve.",
         GOLD,
         ["[HIGH]  dispute      CRT-330575", "[MED]   overpayment  CRT-329181", "[LOW]   refund_open  CRT-331714"]),
    ]

    for i, (title, desc, color, lines) in enumerate(panels):
        x = 0.4 + i * 4.2
        card_rect(s, x, 1.35, 3.95, 5.45, fill=SURF, border=color)
        card_rect(s, x, 1.35, 3.95, 0.06, fill=color, border=color)
        tb(s, title, l=x+0.15, t=1.5, w=3.65, h=0.4,
           size=12, bold=True, color=color)
        tb(s, desc, l=x+0.15, t=1.97, w=3.65, h=2.65,
           size=9.5, color=TEXT)
        card_rect(s, x+0.15, 4.72, 3.65, 0.88, fill=SURF2, border=BORDER)
        ex = "\n".join(lines)
        tb(s, ex, l=x+0.25, t=4.78, w=3.45, h=0.75,
           size=8, color=MUTED, font="Courier New")


def slide_05_payment_journey(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "How a Payment Works",
                subtitle="A citation payment moves through five stages — from 'pay now' to settled funds.")

    steps = [
        ("1", "Cart Created",
         "Customer selects a citation to pay. A cart groups the obligation(s) together for checkout.",
         "Cart:  draft"),
        ("2", "Payment Submitted",
         "Customer enters payment details and confirms. The obligation is locked — preventing duplicate payments.",
         "Cart:  payment_submitted\nObl:   LOCKED"),
        ("3", "Vendor Confirms",
         "Stripe or CheckAlt sends a signal confirming the payment was received.",
         "Cart:  payment_confirmed\nObl:   LOCKED"),
        ("4", "Payment Settles",
         "Funds clear. The ledger is updated, the outstanding balance drops to $0, and the obligation closes.",
         "Cart:  payment_settled\nObl:   CLOSED"),
        ("5", "Receipt Sent",
         "A receipt is automatically delivered. The full payment history is auditable.",
         "Ledger: payment_settled"),
    ]

    for i, (num, title, desc, state) in enumerate(steps):
        x = 0.4 + i * 2.5
        card_rect(s, x, 1.35, 2.35, 5.45, fill=SURF2)
        card_rect(s, x+0.88, 1.5, 0.58, 0.5, fill=GOLD, border=GOLD)
        tb(s, num, l=x+0.88, t=1.52, w=0.58, h=0.42,
           size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)
        tb(s, title, l=x+0.1, t=2.12, w=2.15, h=0.38,
           size=11, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        tb(s, desc, l=x+0.12, t=2.55, w=2.11, h=2.0,
           size=9, color=MUTED)
        card_rect(s, x+0.12, 4.65, 2.11, 0.9, fill=SURF, border=BORDER)
        tb(s, state, l=x+0.2, t=4.7, w=1.95, h=0.8,
           size=8, color=BLUE, font="Courier New")
        if i < 4:
            tb(s, "→", l=x+2.35, t=2.9, w=0.2, h=0.35,
               size=16, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_06_three_concepts(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "The Three Things to Know",
                subtitle="Every citation payment involves three records. Understanding each one helps you answer any question.")

    concepts = [
        ("Obligation",
         "What is owed",
         BLUE,
         [
             "Created when a citation is issued",
             "Tracks: amount owed, amount paid, amount outstanding",
             "Has a status reflecting its current state (OPEN, CLOSED, etc.)",
             "Can be waived or superseded by support staff",
             "Outstanding = amount minus what's been paid via the ledger",
         ],
         ["citation_fee    $40.00   OPEN", "outstanding     $0.00", "status          CLOSED"]),
        ("Cart",
         "The payment session",
         GOLD,
         [
             "Created when a customer starts a checkout",
             "Groups one or more obligations together",
             "Has a status showing where the payment is right now",
             "Links to the vendor transaction (Stripe, CheckAlt)",
             "Refunds are processed against a cart, not an obligation",
         ],
         ["id      CRT-331714", "mode    online (Stripe)", "total   $84.00", "status  payment_settled"]),
        ("Ledger Entry",
         "The permanent record of money moving",
         GREEN,
         [
             "Written automatically when money moves — never modified",
             "Refunds appear as a new negative entry, not a deletion",
             "Every obligation balance is calculated from ledger entries",
             "This is the financial truth — not the cart or obligation fields",
             "If ledger says paid, the obligation is paid — period",
         ],
         ["LDG-1  +$84.00  settled", "  CI-2  OBL-1  +$40", "  CI-1  OBL-9  +$4 ", "LDG-7  -$40.00 refunded"]),
    ]

    for i, (title, subtitle_c, color, bullets, example_lines) in enumerate(concepts):
        x = 0.4 + i * 4.2
        card_rect(s, x, 1.35, 3.95, 5.45, fill=SURF)
        card_rect(s, x, 1.35, 3.95, 0.06, fill=color, border=color)
        tb(s, title, l=x+0.15, t=1.5, w=3.65, h=0.35,
           size=14, bold=True, color=color)
        tb(s, subtitle_c, l=x+0.15, t=1.88, w=3.65, h=0.25,
           size=9.5, italic=True, color=MUTED)
        bullet_lines = [{"text": f"• {b}", "size": 9, "color": TEXT}
                        for b in bullets]
        mtb(s, bullet_lines, l=x+0.15, t=2.18, w=3.65, h=2.0)
        card_rect(s, x+0.12, 4.28, 3.72, 2.3, fill=SURF2, border=BORDER)
        ex_lines = [{"text": ln, "size": 8, "color": MUTED, "font": "Courier New"}
                    for ln in example_lines]
        mtb(s, ex_lines, l=x+0.22, t=4.38, w=3.52, h=2.1,
            default_font="Courier New")


# ─── SECTION 02 ────────────────────────────────────────────────────────────

def slide_08_citation_lookup(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Looking Up a Citation",
                subtitle="Search by citation number to see the full financial picture on one screen.")

    # Browser chrome mock
    card_rect(s, 0.4, 1.28, 12.53, 5.95, fill=SURF2)
    card_rect(s, 0.4, 1.28, 12.53, 0.33, fill=SURF)
    tb(s, "  gringotts.internal / citations / 000322452",
       l=0.5, t=1.3, w=12.2, h=0.26, size=8.5, color=MUTED, font="Courier New")

    # Nav sidebar
    card_rect(s, 0.4, 1.61, 2.0, 5.62, fill=SURF)
    nav = ["Payment Overview", "Obligation History",
           "Cart Detail", "Ledger Entries", "Audit Log", "Refunds"]
    for i, item in enumerate(nav):
        is_active = (i == 0)
        bg  = RGBColor(0x1a, 0x24, 0x38) if is_active else SURF
        col = BLUE if is_active else MUTED
        card_rect(s, 0.45, 1.7 + i*0.52, 1.85, 0.44, fill=bg, border=bg)
        tb(s, item, l=0.55, t=1.73 + i*0.52, w=1.7, h=0.36,
           size=9, color=col)

    # Citation header
    tb(s, "Citation #000322452", l=2.6, t=1.38, w=5, h=0.28,
       size=13, bold=True, color=TEXT)

    # Summary strip
    summary = [
        ("Total Owed",  "$40.00",  TEXT),
        ("Paid",        "$20.00",  GREEN),
        ("Outstanding", "$20.00",  GOLD),
        ("Status",      "PARTIALLY PAID", GOLD),
    ]
    for i, (lbl, val, col) in enumerate(summary):
        x = 2.6 + (i % 2) * 3.7
        y = 1.7 + (i // 2) * 0.44
        tb(s, lbl, l=x, t=y,      w=1.5, h=0.22, size=8,    color=MUTED)
        tb(s, val, l=x, t=y+0.2,  w=2.5, h=0.24, size=11.5, bold=True, color=col)

    tb(s, "Obligations", l=2.6, t=2.65, w=3, h=0.25,
       size=10, bold=True, color=TEXT)
    add_table(s,
              ["label", "owed", "paid", "outstanding", "status"],
              [["citation_fee", "$40.00", "$20.00", "$20.00", "OPEN"]],
              l=2.6, t=2.93, w=9.9, h=0.52,
              font_size=9, col_widths=[0.22, 0.14, 0.14, 0.2, 0.3])

    tb(s, "Carts", l=2.6, t=3.6, w=3, h=0.25,
       size=10, bold=True, color=TEXT)
    add_table(s,
              ["cart id", "mode", "total", "status", "vendor ref"],
              [["CRT-322452-4", "check", "$20.00", "payment_settled", "42/1 (CheckAlt)"]],
              l=2.6, t=3.88, w=9.9, h=0.52,
              font_size=9, col_widths=[0.2, 0.1, 0.1, 0.22, 0.38])

    card_rect(s, 2.6, 4.55, 9.9, 2.0, fill=SURF, border=BORDER)
    tb(s, "How to read this screen", l=2.75, t=4.62, w=9.5, h=0.25,
       size=10, bold=True, color=GOLD)
    notes = [
        "Outstanding = $0.00 and Status = CLOSED means fully paid — no action needed.",
        "Outstanding > $0 and Status = OPEN means the customer still owes money.",
        "Status = LOCKED means a payment is in progress right now — do not take action.",
        "Multiple carts on one citation each represent a separate payment attempt.",
    ]
    mtb(s, [{"text": f"  • {n}", "size": 9, "color": MUTED} for n in notes],
        l=2.75, t=4.9, w=9.5, h=1.5)


def slide_09_obligation_statuses(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Obligation Status Reference",
                subtitle="The obligation status tells you exactly what is happening with a citation's debt.")

    statuses = [
        ("OPEN",       BLUE,   BADGE_OPEN,
         "Unpaid or partially paid",
         "Outstanding balance > $0. The customer has not fully settled this obligation.",
         ["Can be waived (only if no prior payment)",
          "Can be superseded (only if no prior payment)",
          "Will auto-lock when a cart is submitted"]),
        ("LOCKED",     GOLD,   BADGE_LOCKED,
         "Payment in progress",
         "A cart has been submitted and is waiting for vendor confirmation. The obligation is reserved.",
         ["Do not take action while LOCKED",
          "Will auto-close when payment settles",
          "Unlocks automatically if the payment fails"]),
        ("CLOSED",     GREEN,  BADGE_CLOSED,
         "Fully paid",
         "Outstanding balance = $0. Paid in full through the ledger, or remaining balance was waived.",
         ["No action needed on the obligation",
          "A refund can still be issued against the cart",
          "Cannot be waived or superseded"]),
        ("WAIVED",     PURPLE, BADGE_WAIVED,
         "Forgiven — no payment required",
         "The full obligation was waived by support. No payment was received. This is a terminal status.",
         ["Terminal — cannot be changed",
          "Audit log shows who waived it and why",
          "Only applies to fully unpaid obligations"]),
        ("SUPERSEDED", MUTED,  SURF2,
         "Replaced by a new obligation",
         "This obligation was replaced — e.g. the fine amount was corrected. A new obligation was created.",
         ["Terminal — this record is frozen",
          "Find the replacement in the Obligation History tab",
          "Only applies to fully unpaid obligations"]),
        ("DISPUTED",   RED,    BADGE_DISPUTED,
         "Under chargeback dispute",
         "A customer has filed a chargeback. Funds may be withdrawn from the account while the dispute is open.",
         ["Automatically opens a HIGH priority Triage item",
          "Do NOT issue a refund while a dispute is open",
          "May revert to OPEN if the dispute is lost"]),
    ]

    for i, (status, text_col, bg_col, tagline, desc, notes) in enumerate(statuses):
        col = i % 3
        row = i // 3
        x = 0.4  + col * 4.18
        y = 1.35 + row * 2.95
        card_rect(s, x, y, 3.95, 2.75, fill=SURF)
        card_rect(s, x+0.14, y+0.13, 1.35, 0.3, fill=bg_col, border=bg_col)
        tb(s, status, l=x+0.14, t=y+0.14, w=1.35, h=0.26,
           size=8.5, bold=True, color=text_col, align=PP_ALIGN.CENTER)
        tb(s, tagline, l=x+1.6, t=y+0.15, w=2.2, h=0.26,
           size=8.5, italic=True, color=MUTED)
        tb(s, desc, l=x+0.14, t=y+0.52, w=3.67, h=0.7,
           size=9, color=TEXT)
        note_lines = [{"text": f"  • {n}", "size": 8.5, "color": MUTED}
                      for n in notes]
        mtb(s, note_lines, l=x+0.14, t=y+1.26, w=3.67, h=1.38)


def slide_10_cart_statuses(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Cart Status Reference",
                subtitle="The cart status tracks where the payment is — from checkout through settlement and beyond.")

    # Normal journey
    tb(s, "Normal payment journey", l=0.5, t=1.38, w=6, h=0.26,
       size=10, bold=True, color=MUTED)
    normal = [
        ("draft",             MUTED,  SURF2,
         "Cart created. No payment started yet."),
        ("checkout",          BLUE,   BADGE_OPEN,
         "Items added. Customer is on the payment screen."),
        ("payment_submitted", GOLD,   BADGE_LOCKED,
         "Payment sent to vendor. Awaiting confirmation."),
        ("payment_confirmed", BLUE,   BADGE_OPEN,
         "Vendor received payment. Funds not yet cleared."),
        ("payment_settled",   GREEN,  BADGE_CLOSED,
         "Funds cleared. Payment complete."),
    ]
    for i, (status, tcol, bcol, desc) in enumerate(normal):
        x = 0.4 + i * 2.48
        card_rect(s, x, 1.7, 2.3, 1.05, fill=SURF2)
        card_rect(s, x+0.08, 1.78, 2.14, 0.28, fill=bcol, border=bcol)
        tb(s, status, l=x+0.08, t=1.79, w=2.14, h=0.24,
           size=7.5, bold=True, color=tcol, align=PP_ALIGN.CENTER)
        tb(s, desc, l=x+0.1, t=2.1, w=2.1, h=0.55,
           size=8, color=MUTED)
        if i < 4:
            tb(s, "→", l=x+2.3, t=2.05, w=0.22, h=0.3,
               size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Exception states
    tb(s, "Exception states", l=0.5, t=2.95, w=6, h=0.26,
       size=10, bold=True, color=MUTED)
    exceptions = [
        ("abandoned",    MUTED, SURF2,
         "Checkout was not completed within 24 hours. The cart is inactive.",
         "No action needed. The obligation will auto-unlock."),
        ("disputed",     RED,   BADGE_DISPUTED,
         "A chargeback has been filed. Funds may be withdrawn from the account while the dispute is active.",
         "A HIGH priority Triage item is created automatically. Do not issue a refund."),
        ("dispute_lost", RED,   BADGE_DISPUTED,
         "The dispute was resolved in the customer's favor. Funds have been permanently withdrawn.",
         "Terminal status. No further action is possible. An audit record is created."),
    ]
    for i, (status, tcol, bcol, desc, note) in enumerate(exceptions):
        x = 0.4 + i * 4.18
        card_rect(s, x, 3.26, 3.95, 3.72, fill=SURF, border=BORDER)
        card_rect(s, x+0.14, 3.38, 2.0, 0.3, fill=bcol, border=bcol)
        tb(s, status, l=x+0.14, t=3.39, w=2.0, h=0.26,
           size=8.5, bold=True, color=tcol, align=PP_ALIGN.CENTER)
        tb(s, desc, l=x+0.14, t=3.76, w=3.67, h=0.9,
           size=9, color=TEXT)
        card_rect(s, x+0.14, 4.74, 3.67, 1.1, fill=SURF2, border=BORDER)
        tb(s, "What to do:", l=x+0.24, t=4.8, w=3.47, h=0.22,
           size=9, bold=True, color=GOLD)
        tb(s, note, l=x+0.24, t=5.04, w=3.47, h=0.75,
           size=9, color=MUTED)


# ─── SECTION 03 ────────────────────────────────────────────────────────────

def slide_12_action_matrix(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Action Availability — At a Glance",
                subtitle="What you can do depends on the current state of the obligation. Check this first.")

    headers = ["Situation", "Obligation Status", "Refund?", "Full Waive?", "Partial Waive?", "Supersede?"]
    rows = [
        ["Not yet paid at all",    "OPEN  (paid = $0)",    "✕ No settled cart", "✓ Yes",                "✕ Nothing paid yet",  "✓ Yes"],
        ["Partially paid",         "OPEN  (paid > $0)",    "✓ Via settled cart", "✕ Use partial waive", "✓ Yes",               "✕ Has payment"],
        ["Fully paid",             "CLOSED",               "✓ Via settled cart", "✕ Already closed",    "✕ Already closed",    "✕ Already closed"],
        ["Payment in progress",    "LOCKED",               "✕ Not settled yet",  "✕ Wait for payment",  "✕ Wait for payment",  "✕ Wait for payment"],
        ["Waived / Superseded",    "WAIVED / SUPERSEDED",  "✕ Terminal",         "✕ Terminal",          "✕ Terminal",          "✕ Terminal"],
    ]
    add_table(s, headers, rows, l=0.4, t=1.35, w=12.53, h=3.6,
              font_size=9, col_widths=[0.22, 0.18, 0.15, 0.15, 0.15, 0.15])

    card_rect(s, 0.4, 5.12, 12.53, 1.55, fill=RGBColor(0x1a, 0x16, 0x00), border=GOLD)
    tb(s, "Rules to remember:", l=0.6, t=5.18, w=5, h=0.26,
       size=10, bold=True, color=GOLD)
    rules = [
        "Waivers and supersedes require the obligation to have NO prior payments (paid = $0).",
        "Partial waive is only for obligations that are partially paid (paid > $0, outstanding > $0).",
        "Refunds are against a cart, not an obligation. The cart must be in payment_settled status.",
        "When in doubt about a LOCKED obligation, wait — it will resolve on its own within minutes.",
    ]
    mtb(s, [{"text": f"  • {r}", "size": 9, "color": MUTED} for r in rules],
        l=0.6, t=5.45, w=12.2, h=1.1)


def slide_13_refund(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Issuing a Refund",
                subtitle="Refunds reverse a payment. The cart must be in payment_settled status.")

    # Left column
    card_rect(s, 0.4, 1.35, 5.85, 5.45, fill=SURF)
    tb(s, "When to issue a refund", l=0.55, t=1.42, w=5.55, h=0.28,
       size=11, bold=True, color=GOLD)
    when = [
        "Customer was charged in error or double-charged",
        "Payment was made but the citation was dismissed",
        "Customer paid the wrong citation",
        "Court ordered a reversal of the payment",
    ]
    for i, w_text in enumerate(when):
        tb(s, f"• {w_text}", l=0.65, t=1.74 + i*0.3, w=5.5, h=0.27,
           size=9.5, color=MUTED)

    tb(s, "Steps", l=0.55, t=3.0, w=5.55, h=0.28,
       size=11, bold=True, color=TEXT)
    steps = [
        ("Go to Cart Detail",
         "Find the citation, click Cart Detail. Confirm the cart status is payment_settled."),
        ("Select items and enter amount",
         "Check the cart items to include. Enter the refund amount and a reason."),
        ("Submit the refund",
         "Click Initiate Refund. The system processes this asynchronously — you'll see the result in minutes."),
        ("Ledger is updated",
         "A new negative entry is appended to the ledger. The original entry is never changed."),
        ("Obligation balance recalculates",
         "Outstanding balance updates automatically. The obligation may reopen if fully reversed."),
    ]
    for i, (title, desc) in enumerate(steps):
        y = 3.32 + i * 0.57
        card_rect(s, 0.55, y+0.03, 0.36, 0.34, fill=GOLD, border=GOLD)
        tb(s, str(i+1), l=0.55, t=y+0.04, w=0.36, h=0.3,
           size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)
        tb(s, title, l=1.0, t=y+0.02, w=4.75, h=0.22,
           size=9.5, bold=True, color=TEXT)
        tb(s, desc,  l=1.0, t=y+0.24, w=4.75, h=0.3,
           size=8.5, color=MUTED)

    # Right column
    card_rect(s, 6.5, 1.35, 6.43, 5.45, fill=SURF)
    tb(s, "What you'll see in the Ledger Report", l=6.65, t=1.42, w=6.1, h=0.28,
       size=11, bold=True, color=BLUE)

    tb(s, "Before the refund:", l=6.65, t=1.8, w=6.1, h=0.25,
       size=9.5, bold=True, color=TEXT)
    before_hdrs = ["id",    "type",            "amount",  "stage"]
    before_rows = [
        ["LDG-1", "cart_payment",     "+$84.00", "payment_settled"],
        ["LDG-2", "stripe_fee",       "-$1.55",  "payment_settled"],
    ]
    add_table(s, before_hdrs, before_rows,
              l=6.65, t=2.08, w=6.1, h=0.88,
              font_size=8.5, col_widths=[0.12, 0.35, 0.2, 0.33])

    tb(s, "After issuing a $40 refund on OBL-1:", l=6.65, t=3.1, w=6.1, h=0.25,
       size=9.5, bold=True, color=TEXT)
    after_hdrs = ["id",    "type",              "amount",  "stage"]
    after_rows = [
        ["LDG-1", "cart_payment",     "+$84.00", "payment_settled"],
        ["LDG-2", "stripe_fee",       "-$1.55",  "payment_settled"],
        ["LDG-7", "payment_refunded", "-$40.00", "payment_refunded"],
    ]
    add_table(s, after_hdrs, after_rows,
              l=6.65, t=3.38, w=6.1, h=1.05,
              font_size=8.5, col_widths=[0.12, 0.35, 0.2, 0.33])

    card_rect(s, 6.65, 4.55, 6.1, 0.95, fill=SURF2, border=GREEN)
    tb(s, "OBL-1 balance after refund:", l=6.75, t=4.6, w=5.9, h=0.25,
       size=9.5, bold=True, color=TEXT)
    tb(s, "Paid:           $40 − $40 = $0.00\n"
          "Outstanding:    $40 − $0  = $40.00\n"
          "Status:         OPEN  (reopened)",
       l=6.75, t=4.87, w=5.9, h=0.58,
       size=8.5, color=GREEN, font="Courier New")

    card_rect(s, 6.65, 5.62, 6.1, 0.95, fill=RGBColor(0x1a, 0x16, 0x00), border=GOLD)
    tb(s, "The ledger is never modified. Refunds are always a new row appended to the bottom.",
       l=6.75, t=5.75, w=5.9, h=0.65,
       size=9.5, italic=True, color=GOLD)


def slide_14_waive(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Waiving an Obligation",
                subtitle="Forgive a financial obligation so that no payment is required.")

    types = [
        ("Full Waive",
         "Forgive the entire outstanding balance when no payment has been made.",
         PURPLE,
         ["Obligation status = OPEN", "No prior payments (paid = $0)", "Outstanding balance > $0"],
         ["Status → WAIVED  (terminal, cannot be changed)",
          "Outstanding balance → $0",
          "Waived amount and reason recorded in audit log"],
         "Court grants a full hardship exemption.\nCitation is dismissed — no payment required."),
        ("Partial Waive",
         "Forgive the remaining balance after the customer has already made a partial payment.",
         GOLD,
         ["Obligation status = OPEN", "Has prior payments (paid > $0)", "Outstanding balance > $0"],
         ["Status → CLOSED  (obligation is now settled)",
          "Remaining outstanding balance → $0",
          "Waived amount and reason recorded in audit log"],
         "Customer paid $500 of a $1,000 fine.\nCourt agrees to forgive the remaining $500."),
    ]

    for i, (title, desc, color, preconditions, effects, use_case) in enumerate(types):
        x = 0.4 + i * 6.3
        card_rect(s, x, 1.35, 6.05, 5.45, fill=SURF)
        card_rect(s, x, 1.35, 6.05, 0.06, fill=color, border=color)
        tb(s, title, l=x+0.15, t=1.5, w=5.75, h=0.35,
           size=15, bold=True, color=color)
        tb(s, desc, l=x+0.15, t=1.9, w=5.75, h=0.45,
           size=9.5, color=MUTED)

        tb(s, "Required conditions:", l=x+0.15, t=2.43, w=5.75, h=0.26,
           size=10, bold=True, color=TEXT)
        for j, pc in enumerate(preconditions):
            tb(s, f"  • {pc}", l=x+0.15, t=2.72 + j*0.28, w=5.75, h=0.26,
               size=9, color=MUTED)

        tb(s, "What happens:", l=x+0.15, t=3.6, w=5.75, h=0.26,
           size=10, bold=True, color=TEXT)
        for j, eff in enumerate(effects):
            tb(s, f"  • {eff}", l=x+0.15, t=3.88 + j*0.28, w=5.75, h=0.26,
               size=9, color=TEXT)

        card_rect(s, x+0.15, 4.82, 5.75, 1.72, fill=SURF2, border=BORDER)
        tb(s, "Example use case:", l=x+0.25, t=4.88, w=5.55, h=0.24,
           size=9.5, bold=True, color=GOLD)
        tb(s, use_case, l=x+0.25, t=5.14, w=5.55, h=1.3,
           size=9.5, color=MUTED)


def slide_15_supersede(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Superseding an Obligation",
                subtitle="Replace an existing obligation with a corrected amount — for example, when a court revises a fine.")

    card_rect(s, 0.4, 1.35, 5.75, 5.45, fill=SURF)
    tb(s, "What supersede does", l=0.55, t=1.42, w=5.45, h=0.28,
       size=12, bold=True, color=GOLD)
    tb(s, "Use supersede when the fine amount must be changed. The original obligation is "
          "frozen and marked SUPERSEDED, and a brand-new obligation is created at the corrected amount.",
       l=0.55, t=1.78, w=5.45, h=0.72, size=9.5, color=TEXT)

    tb(s, "Required conditions:", l=0.55, t=2.58, w=5.45, h=0.26,
       size=10, bold=True, color=TEXT)
    preconds = ["Obligation status = OPEN",
                "No prior payments have been made (paid = $0)",
                "You have the corrected amount ready"]
    for j, pc in enumerate(preconds):
        tb(s, f"  • {pc}", l=0.55, t=2.86 + j*0.28, w=5.45, h=0.26,
           size=9, color=MUTED)

    tb(s, "What happens:", l=0.55, t=3.72, w=5.45, h=0.26,
       size=10, bold=True, color=TEXT)
    effects = [
        "Original obligation status → SUPERSEDED  (terminal)",
        "A new obligation is created with the corrected amount",
        "New obligation status = OPEN  (customer still owes)",
        "Both old and new obligations appear in Obligation History",
        "Audit log records the action with full before/after snapshots",
    ]
    for j, eff in enumerate(effects):
        tb(s, f"  • {eff}", l=0.55, t=4.0 + j*0.28, w=5.45, h=0.26,
           size=9, color=TEXT)

    card_rect(s, 0.55, 5.58, 5.45, 0.98, fill=SURF2, border=RED)
    tb(s, "Cannot supersede if:", l=0.65, t=5.64, w=5.25, h=0.24,
       size=9.5, bold=True, color=RED)
    tb(s, "Any payment has been applied to the obligation (paid > $0). "
          "Contact engineering for a manual adjustment in that case.",
       l=0.65, t=5.9, w=5.25, h=0.58, size=9, color=MUTED)

    # Before / after visual
    card_rect(s, 6.45, 1.35, 6.43, 5.45, fill=SURF)
    tb(s, "Before supersede:", l=6.6, t=1.42, w=6.1, h=0.28,
       size=11, bold=True, color=TEXT)
    before = [("ID",          "OBL-5"),
              ("Label",       "citation_fee"),
              ("Amount",      "$40.00"),
              ("Paid",        "$0.00"),
              ("Outstanding", "$40.00"),
              ("Status",      "OPEN")]
    for j, (lbl, val) in enumerate(before):
        y = 1.78 + j * 0.3
        tb(s, lbl,  l=6.6, t=y, w=2.0, h=0.26, size=9, color=MUTED)
        col = BLUE if val == "OPEN" else GOLD if val == "$40.00" and lbl == "Outstanding" else TEXT
        tb(s, val,  l=8.7, t=y, w=2.5, h=0.26, size=9, bold=True, color=col)

    card_rect(s, 6.6, 3.6, 6.1, 0.38, fill=RGBColor(0x1a, 0x16, 0x00), border=GOLD)
    tb(s, "  ↓   Supersede — new amount: $25.00",
       l=6.6, t=3.63, w=6.1, h=0.3, size=10, bold=True, color=GOLD)

    tb(s, "After supersede:", l=6.6, t=4.08, w=6.1, h=0.28,
       size=11, bold=True, color=TEXT)

    card_rect(s, 6.6, 4.4, 2.9, 2.08, fill=SURF2, border=MUTED)
    tb(s, "OLD — SUPERSEDED", l=6.7, t=4.46, w=2.7, h=0.24,
       size=8.5, bold=True, color=MUTED)
    tb(s, "OBL-5  $40.00\nStatus: SUPERSEDED\nFrozen. No further action.",
       l=6.7, t=4.72, w=2.7, h=0.85, size=8.5, color=MUTED, font="Courier New")

    card_rect(s, 9.75, 4.4, 2.95, 2.08, fill=SURF2, border=GREEN)
    tb(s, "NEW — OPEN", l=9.85, t=4.46, w=2.75, h=0.24,
       size=8.5, bold=True, color=GREEN)
    tb(s, "OBL-6  $25.00\nStatus: OPEN\nCustomer owes $25.00.",
       l=9.85, t=4.72, w=2.75, h=0.85, size=8.5, color=GREEN, font="Courier New")


# ─── SECTION 04 ────────────────────────────────────────────────────────────

def slide_17_tracing_complaint(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Tracing a Complaint",
                subtitle="Every action in Gringotts is logged with a correlation ID — use it to follow the chain of events.")

    card_rect(s, 0.4, 1.35, 5.75, 5.45, fill=SURF)
    tb(s, "Step-by-step investigation", l=0.55, t=1.42, w=5.45, h=0.28,
       size=11, bold=True, color=TEXT)

    trace_steps = [
        ("Get the citation number",
         "Ask the customer for the citation number. Every investigation starts here."),
        ("Check the obligation status",
         "Is it OPEN, CLOSED, LOCKED? This tells you immediately whether money is still owed."),
        ("Find the relevant cart",
         "Go to Cart Detail. Each cart is a payment attempt. Identify the one the customer is asking about."),
        ("Open the Audit Log",
         "The audit log shows every status change in order. Look for the event chain around the time of the complaint."),
        ("Search by correlation ID",
         "Copy the correlation_id from any audit log entry. Search across all logs with this ID to see exactly what happened — webhook received, ledger written, obligation updated, etc."),
    ]
    for i, (title, desc) in enumerate(trace_steps):
        y = 1.78 + i * 1.0
        card_rect(s, 0.55, y+0.05, 0.36, 0.34, fill=PURPLE, border=PURPLE)
        tb(s, str(i+1), l=0.55, t=y+0.06, w=0.36, h=0.3,
           size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(s, title, l=1.0, t=y+0.04, w=4.78, h=0.24,
           size=10, bold=True, color=TEXT)
        tb(s, desc, l=1.0, t=y+0.3, w=4.78, h=0.62,
           size=8.5, color=MUTED)

    card_rect(s, 6.4, 1.35, 6.43, 5.45, fill=SURF)
    tb(s, "Sample audit log for Citation #000331714", l=6.55, t=1.42, w=6.1, h=0.28,
       size=11, bold=True, color=TEXT)

    audit_hdrs = ["event", "entity", "change", "actor", "correlation"]
    audit_rows = [
        ["obligation.locked",            "OBL-1", "OPEN → LOCKED",    "system",  "cart_payment:1"],
        ["obligation.payment_confirmed", "OBL-1", "LOCKED → LOCKED",  "system",  "cart_payment:1"],
        ["obligation.payment_settled",   "OBL-1", "LOCKED → CLOSED",  "system",  "cart_payment:1"],
        ["cart.payment_settled",         "CRT-1", "confirmed→settled","system",  "cart_payment:1"],
        ["obligation.payment_refunded",  "OBL-1", "CLOSED → OPEN",    "support", "refund:7"],
    ]
    add_table(s, audit_hdrs, audit_rows,
              l=6.55, t=1.78, w=6.1, h=1.88,
              font_size=8.5, col_widths=[0.32, 0.12, 0.24, 0.12, 0.2])

    tb(s, "What the correlation ID tells you:", l=6.55, t=3.78, w=6.1, h=0.26,
       size=10, bold=True, color=GOLD)
    tb(s, "cart_payment:1",
       l=6.55, t=4.08, w=3.5, h=0.3,
       size=12, bold=True, color=BLUE, font="Courier New")
    explain = [
        "All five events above share the same correlation ID — they are all part of the same payment.",
        "The final entry (refund:7) is a different correlation ID, meaning it was a separate action.",
        "If a customer says 'I paid but my obligation is OPEN' — find their cart, get the correlation_id, search the audit log. If payment_settled is missing, the vendor event may not have arrived yet.",
    ]
    mtb(s, [{"text": f"• {e}", "size": 9, "color": MUTED} for e in explain],
        l=6.55, t=4.42, w=6.1, h=2.25)


def slide_18_common_scenarios(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Common Complaint Scenarios",
                subtitle="The most frequent issues you'll see and exactly how to handle each one.")

    scenarios = [
        ('"I paid but it still shows I owe money."',
         GOLD,
         ["Look up the citation — check the obligation outstanding balance.",
          "Find the cart. Is it payment_settled? If yes, check the obligation activity log.",
          "If payment settled but obligation still OPEN: look for a refund that reopened it.",
          "If cart is in payment_submitted for 24+ hours: create a Stuck Payment triage item.",
          "If no cart exists at all: the payment may not have been processed — ask for receipt."]),
        ('"I was charged twice for the same citation."',
         RED,
         ["Look up the citation. You should see two carts in payment_settled.",
          "Confirm the obligation shows an overpayment (paid > owed).",
          "Issue a refund on the duplicate cart via the Refund tab.",
          "If the second cart is in payment_submitted (not settled): wait — it may auto-fail within 24hrs.",
          "If both are settled and it's been 24hrs: issue the refund immediately."]),
        ('"The system says it\'s locked and I can\'t pay."',
         BLUE,
         ["LOCKED means a payment is actively in progress on another cart.",
          "Check the cart list — find the cart in payment_submitted.",
          "If it's been less than 1 hour: tell the customer to wait for confirmation.",
          "If it's been 24+ hours with no update: create a Stuck Payment triage item.",
          "The obligation unlocks automatically once the cart is resolved."]),
        ('"My bank reversed the charge / chargeback."',
         RED,
         ["Gringotts auto-creates a HIGH priority Triage item when this happens.",
          "The cart status changes to disputed; funds are temporarily withdrawn.",
          "Do NOT issue a refund while a dispute is open — it creates a double-reversal.",
          "Gather evidence (receipt, payment date, obligation details) for the dispute.",
          "If dispute is won: cart returns to payment_settled. If lost: dispute_lost (terminal)."]),
    ]

    for i, (question, color, steps) in enumerate(scenarios):
        row = i // 2
        col = i % 2
        x = 0.4  + col * 6.45
        y = 1.35 + row * 3.05
        card_rect(s, x, y, 6.2, 2.88, fill=SURF, border=BORDER)
        card_rect(s, x, y, 6.2, 0.06, fill=color, border=color)
        tb(s, question, l=x+0.15, t=y+0.13, w=5.9, h=0.38,
           size=10, bold=True, color=TEXT, italic=True)
        step_lines = [{"text": f"  {j+1}. {step}", "size": 8.5, "color": MUTED}
                      for j, step in enumerate(steps)]
        mtb(s, step_lines, l=x+0.15, t=y+0.56, w=5.9, h=2.18)


# ─── SECTION 05 ────────────────────────────────────────────────────────────

def slide_20_triage_queue(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "The Triage Queue",
                subtitle="Gringotts automatically detects payment anomalies and surfaces them here for you to investigate and resolve.")

    tb(s, "Current queue (4 open items):", l=0.4, t=1.38, w=12.53, h=0.26,
       size=10, bold=True, color=TEXT)

    triage_hdrs = ["priority", "type", "citation", "cart", "detail", "assigned", "status"]
    triage_rows = [
        ["HIGH", "dispute",         "000330575", "Cart 2",    "Chargeback filed $84.00",                  "sarah.k",  "in_progress"],
        ["MED",  "overpayment",     "000329181", "Carts 1,3", "Paid by Stripe + CheckAlt ($40 over)",     "—",        "open"],
        ["MED",  "partial_payment", "000322452", "Cart 4",    "Check paid $20, still owes $20",           "james.r",  "in_progress"],
        ["LOW",  "refund_open",     "000331714", "Cart 1",    "Stripe refund issued, OBL-1 still OPEN",   "—",        "open"],
    ]
    add_table(s, triage_hdrs, triage_rows,
              l=0.4, t=1.67, w=12.53, h=1.52,
              font_size=9,
              col_widths=[0.08, 0.14, 0.13, 0.1, 0.3, 0.12, 0.13])

    tb(s, "Auto-detection rules — how items enter the queue:", l=0.4, t=3.35, w=12.53, h=0.26,
       size=10, bold=True, color=TEXT)

    type_hdrs = ["Type", "How it's detected", "Priority", "What to do"]
    type_rows = [
        ["dispute",
         "Vendor sends a payment_dispute_funds_withdrawn signal",
         "HIGH",
         "Gather evidence. Do NOT issue a refund. Track the dispute window deadline with the vendor."],
        ["stuck_payment",
         "Cart stays in payment_submitted for more than 24 hours",
         "MED",
         "Check with Stripe or CheckAlt for the payment intent status. Cancel and let the customer retry."],
        ["refund_failed",
         "Refund processor has returned an error 3 or more times",
         "MED",
         "Issue the refund manually through the vendor portal. Mark the Triage item resolved."],
        ["balance_mismatch",
         "Obligation outstanding > $0 but linked cart is payment_settled",
         "MED",
         "Inspect the ledger entries for the cart. Check for a missing event. Escalate to engineering if stuck."],
        ["vendor_event_failed",
         "A vendor webhook could not be processed after retries",
         "LOW",
         "Inspect the error in the event log. Fix bad data or missing cart. Ask engineering to replay the event."],
    ]
    add_table(s, type_hdrs, type_rows,
              l=0.4, t=3.65, w=12.53, h=2.85,
              font_size=8.5,
              col_widths=[0.14, 0.3, 0.09, 0.47])


def slide_21_quick_reference(prs):
    s = blank_slide(prs)
    set_bg(s)
    slide_title(s, "Quick Reference Card",
                subtitle="Keep this handy. Everything you need to answer common questions at a glance.")

    # ── Obligation statuses ─────────────────────────────────────────────────
    card_rect(s, 0.4, 1.35, 3.95, 5.45, fill=SURF)
    tb(s, "Obligation Status", l=0.55, t=1.42, w=3.65, h=0.26,
       size=11, bold=True, color=TEXT)

    obl_ref = [
        ("OPEN",       BLUE,   BADGE_OPEN,    "Unpaid or partially paid"),
        ("LOCKED",     GOLD,   BADGE_LOCKED,  "Payment in progress"),
        ("CLOSED",     GREEN,  BADGE_CLOSED,  "Fully paid"),
        ("WAIVED",     PURPLE, BADGE_WAIVED,  "Forgiven — no payment"),
        ("SUPERSEDED", MUTED,  SURF2,         "Replaced by new obligation"),
        ("DISPUTED",   RED,    BADGE_DISPUTED,"Under chargeback dispute"),
    ]
    for i, (status, col, bg, desc) in enumerate(obl_ref):
        y = 1.82 + i * 0.54
        card_rect(s, 0.55, y+0.05, 1.3, 0.28, fill=bg, border=bg)
        tb(s, status, l=0.55, t=y+0.06, w=1.3, h=0.24,
           size=8, bold=True, color=col, align=PP_ALIGN.CENTER)
        tb(s, desc, l=1.95, t=y+0.07, w=2.25, h=0.24,
           size=9, color=MUTED)

    # ── Action preconditions ────────────────────────────────────────────────
    card_rect(s, 4.6, 1.35, 4.2, 5.45, fill=SURF)
    tb(s, "Action Preconditions", l=4.75, t=1.42, w=3.9, h=0.26,
       size=11, bold=True, color=TEXT)

    action_ref = [
        ("Refund",        "Cart = payment_settled",        GREEN),
        ("Full Waive",    "OBL = OPEN  &  paid = $0",      PURPLE),
        ("Partial Waive", "OBL = OPEN  &  paid > $0",      GOLD),
        ("Supersede",     "OBL = OPEN  &  paid = $0",      BLUE),
        ("No action",     "CLOSED / LOCKED / any terminal", MUTED),
    ]
    for i, (action, condition, col) in enumerate(action_ref):
        y = 1.82 + i * 0.68
        tb(s, action,    l=4.75, t=y,      w=1.7, h=0.28, size=10, bold=True, color=col)
        tb(s, condition, l=4.75, t=y+0.3,  w=3.9, h=0.28, size=8.5, color=MUTED,
           font="Courier New")
        if i < 4:
            divider = s.shapes.add_shape(
                1, Inches(4.75), Inches(y+0.6), Inches(3.9), Inches(0.01))
            divider.fill.solid(); divider.fill.fore_color.rgb = BORDER
            divider.line.fill.background()

    # ── Triage items ────────────────────────────────────────────────────────
    card_rect(s, 9.05, 1.35, 3.88, 5.45, fill=SURF)
    tb(s, "Triage Item Types", l=9.2, t=1.42, w=3.58, h=0.26,
       size=11, bold=True, color=TEXT)

    triage_ref = [
        ("dispute",             RED,  "HIGH", "Chargeback filed against payment"),
        ("stuck_payment",       GOLD, "MED",  "24h+ in payment_submitted"),
        ("refund_failed",       GOLD, "MED",  "3+ refund processing errors"),
        ("balance_mismatch",    GOLD, "MED",  "Settled cart, OBL still open"),
        ("vendor_event_failed", MUTED,"LOW",  "Webhook processing error"),
    ]
    for i, (ttype, col, pri, desc) in enumerate(triage_ref):
        y = 1.82 + i * 0.72
        pri_col = RED if pri == "HIGH" else (GOLD if pri == "MED" else MUTED)
        tb(s, ttype, l=9.2, t=y, w=2.35, h=0.26,
           size=9.5, bold=True, color=col)
        card_rect(s, 11.65, y+0.02, 0.88, 0.24, fill=SURF2, border=pri_col)
        tb(s, pri, l=11.65, t=y+0.03, w=0.88, h=0.22,
           size=8, bold=True, color=pri_col, align=PP_ALIGN.CENTER)
        tb(s, desc, l=9.2, t=y+0.3, w=3.5, h=0.3,
           size=8.5, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    print("Building Gringotts Support Guide...")

    slide_01_title(prs)                                      ; print("  1/21 Title")
    slide_02_agenda(prs)                                     ; print("  2/21 Agenda")

    section_divider(prs, "01", "How Gringotts\nWorks",
        "What the system does, how a payment travels from citation to settled,\n"
        "and the three key concepts every support person needs to know.", BLUE)  ; print("  3/21 Section: How It Works")
    slide_04_what_is_gringotts(prs)                          ; print("  4/21 What is Gringotts?")
    slide_05_payment_journey(prs)                            ; print("  5/21 Payment Journey")
    slide_06_three_concepts(prs)                             ; print("  6/21 Three Key Concepts")

    section_divider(prs, "02", "Reading a\nCitation's State",
        "How to look up a citation, interpret obligation and cart statuses,\n"
        "and understand what the numbers mean.", GREEN)       ; print("  7/21 Section: Reading State")
    slide_08_citation_lookup(prs)                            ; print("  8/21 Citation Lookup")
    slide_09_obligation_statuses(prs)                        ; print("  9/21 Obligation Statuses")
    slide_10_cart_statuses(prs)                              ; print(" 10/21 Cart Statuses")

    section_divider(prs, "03", "Actions You\nCan Take",
        "When and how to issue a refund, waive an obligation, or supersede\n"
        "an obligation — and what you cannot change.", GOLD)  ; print(" 11/21 Section: Actions")
    slide_12_action_matrix(prs)                              ; print(" 12/21 Action Matrix")
    slide_13_refund(prs)                                     ; print(" 13/21 Issuing a Refund")
    slide_14_waive(prs)                                      ; print(" 14/21 Waiving an Obligation")
    slide_15_supersede(prs)                                  ; print(" 15/21 Superseding an Obligation")

    section_divider(prs, "04", "Investigating\nComplaints",
        "How to trace a payment issue from complaint to root cause\n"
        "using audit logs and correlation IDs.", PURPLE)      ; print(" 16/21 Section: Investigating")
    slide_17_tracing_complaint(prs)                          ; print(" 17/21 Tracing a Complaint")
    slide_18_common_scenarios(prs)                           ; print(" 18/21 Common Scenarios")

    section_divider(prs, "05", "The Triage\nQueue",
        "Your daily work queue: item types, priorities, auto-detection,\n"
        "and how to resolve each one.", RED)                  ; print(" 19/21 Section: Triage")
    slide_20_triage_queue(prs)                               ; print(" 20/21 Triage Queue")
    slide_21_quick_reference(prs)                            ; print(" 21/21 Quick Reference Card")

    out = "gringotts_support_guide.pptx"
    prs.save(out)
    print(f"\nSaved → /Users/obvio/Documents/code_repos/gringotts/{out}")


if __name__ == "__main__":
    main()
